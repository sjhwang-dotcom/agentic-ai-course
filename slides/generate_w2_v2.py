#!/usr/bin/env python3
"""Generate W2 Agentic Architecture Patterns lecture slides (~78 slides).
Redesigned with hierarchical structure, transitions, custom diagrams, and HQ figures."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Load reference template and strip slides ───
ref_path = "/Users/sungju/Downloads/Agentic AI - Intro.pptx"
prs = Presentation(ref_path)

while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# ─── Style constants ───
FONT_TITLE = "Arial"
FONT_BODY = "Inter"
FONT_BODY_LIGHT = "Inter Light"
COLOR_BODY = RGBColor(0x43, 0x43, 0x43)
COLOR_REF = RGBColor(0x66, 0x66, 0x66)
COLOR_ACCENT = RGBColor(0x42, 0x85, 0xF4)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BLACK = RGBColor(0x00, 0x00, 0x00)
COLOR_DK2 = RGBColor(0x59, 0x59, 0x59)

# Section colors
COLOR_BLUE = RGBColor(0x42, 0x85, 0xF4)
COLOR_GREEN = RGBColor(0x0F, 0x9D, 0x58)
COLOR_ORANGE = RGBColor(0xF4, 0xB4, 0x00)
COLOR_PURPLE = RGBColor(0xAB, 0x47, 0xBC)
COLOR_RED = RGBColor(0xDB, 0x44, 0x37)
COLOR_TEAL = RGBColor(0x00, 0x97, 0xA7)
COLOR_LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
COLOR_LIGHT_BLUE_BG = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_LIGHT_GREEN_BG = RGBColor(0xE6, 0xF4, 0xEA)
COLOR_LIGHT_ORANGE_BG = RGBColor(0xFE, 0xF7, 0xE0)

FIG_DIR = "figures/hq"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


# ═══════════════════════════════════════════════════════════════════
#  HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════════════

def add_title_slide(title_text, subtitle_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    ph = slide.placeholders[0]
    ph.text = title_text
    for p in ph.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = COLOR_BLACK
    if subtitle_text and 1 in slide.placeholders:
        sub = slide.placeholders[1]
        sub.text = subtitle_text
        for p in sub.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = COLOR_DK2
    return slide


def add_section_slide(title_text, subtitle_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    ph = slide.placeholders[0]
    ph.text = title_text
    for p in ph.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = COLOR_BLACK
    if subtitle_text and 1 in slide.placeholders:
        sub = slide.placeholders[1]
        sub.text = subtitle_text
        for p in sub.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = COLOR_DK2
    return slide


def _add_run(para, text, size=Pt(11), bold=False, color=None, font_name=None):
    run = para.add_run()
    run.text = text
    run.font.size = size
    run.font.name = font_name or (FONT_BODY if bold else FONT_BODY_LIGHT)
    run.font.bold = bold
    run.font.color.rgb = color or COLOR_BODY
    return run


def add_content_slide(title, items, ref="", fig=None, fig_pos=None,
                      body_top=None, body_width=None, bullet_size=Pt(11)):
    """items: list of (text, bold, indent) or str."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    ph = slide.placeholders[0]
    ph.text = title
    for p in ph.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = COLOR_BLACK

    fig_path = os.path.join(FIG_DIR, fig) if fig else None
    has_fig = fig_path and os.path.exists(fig_path)

    if has_fig:
        tw = body_width or Inches(4.8)
        tl, tt = Inches(0.4), body_top or Inches(1.2)
        th = Inches(3.6)
        if fig_pos:
            fl, ft, fw, fh = fig_pos
        else:
            fl, ft, fw, fh = Inches(5.4), Inches(1.2), Inches(4.2), Inches(3.2)
        try:
            slide.shapes.add_picture(fig_path, fl, ft, fw, fh)
        except Exception:
            pass
    else:
        tw = body_width or Inches(9.0)
        tl, tt = Inches(0.4), body_top or Inches(1.2)
        th = Inches(3.8)

    txBox = slide.shapes.add_textbox(tl, tt, tw, th)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, bold, indent = item
        else:
            text, bold, indent = item, False, 0
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _add_run(para, text, size=bullet_size, bold=bold)
        para.level = indent
        para.space_after = Pt(4)

    if ref:
        rb = slide.shapes.add_textbox(Inches(0.4), Inches(5.05), Inches(9.0), Inches(0.4))
        rtf = rb.text_frame
        rtf.word_wrap = True
        _add_run(rtf.paragraphs[0], ref, size=Pt(7), color=COLOR_REF, font_name=FONT_BODY_LIGHT)

    return slide


def add_figure_slide(title, fig, ref="", caption="", fig_top=None, fig_h=None):
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    ph = slide.placeholders[0]
    ph.text = title
    for p in ph.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = COLOR_BLACK

    fig_path = os.path.join(FIG_DIR, fig)
    if os.path.exists(fig_path):
        from PIL import Image
        img = Image.open(fig_path)
        aspect = img.width / img.height
        max_w = Inches(8.5)
        max_h = fig_h or Inches(3.5)
        if aspect > (max_w / max_h):
            w = max_w; h = int(w / aspect)
        else:
            h = max_h; w = int(h * aspect)
        left = int((SLIDE_W - w) / 2)
        top = fig_top or Inches(1.3)
        try:
            slide.shapes.add_picture(fig_path, left, top, w, h)
        except Exception:
            pass

    if caption:
        cb = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9.0), Inches(0.3))
        _add_run(cb.text_frame.paragraphs[0], caption, size=Pt(9), color=COLOR_DK2,
                 font_name=FONT_BODY_LIGHT)

    if ref:
        rb = slide.shapes.add_textbox(Inches(0.4), Inches(5.05), Inches(9.0), Inches(0.4))
        _add_run(rb.text_frame.paragraphs[0], ref, size=Pt(7), color=COLOR_REF,
                 font_name=FONT_BODY_LIGHT)
    return slide


def add_transition_box(slide, text, color=COLOR_ACCENT):
    """Add a colored transition box at the bottom of a slide."""
    left, top = Inches(0.4), Inches(4.45)
    width, height = Inches(9.2), Inches(0.55)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Adjust corner radius
    shape.adjustments[0] = 0.1
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    _add_run(para, text, size=Pt(9), color=COLOR_WHITE, bold=True, font_name=FONT_BODY)


def draw_box(slide, left, top, w, h, text, fill, text_color=COLOR_WHITE,
             font_size=Pt(9), bold=True, align=PP_ALIGN.CENTER):
    """Draw a labeled rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.15
    except Exception:
        pass
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    para = tf.paragraphs[0]
    para.alignment = align
    _add_run(para, text, size=font_size, bold=bold, color=text_color, font_name=FONT_BODY)
    return shape


def draw_arrow_right(slide, left, top, width, color=COLOR_DK2):
    """Draw a right-pointing arrow."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def draw_label(slide, left, top, w, h, text, size=Pt(9), color=COLOR_BODY, bold=False,
               align=PP_ALIGN.CENTER):
    """Add a text label (no background)."""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    _add_run(para, text, size=size, color=color, bold=bold, font_name=FONT_BODY if bold else FONT_BODY_LIGHT)
    return tb


# ═══════════════════════════════════════════════════════════════════
#  PART 0: OPENING (Slides 1-7)
# ═══════════════════════════════════════════════════════════════════

def generate_part0():
    # --- Slide 1: Title ---
    add_title_slide("Agentic Architecture Patterns",
                    "AI 89900 — Agentic AI  |  W2  |  KAIST Spring 2026")

    # --- Slide 2: The Central Question ---
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.placeholders[0].text = "The Central Question"
    for p in slide.placeholders[0].text_frame.paragraphs:
        for r in p.runs: r.font.color.rgb = COLOR_BLACK

    # Main question
    draw_label(slide, Inches(0.5), Inches(1.1), Inches(9.0), Inches(0.6),
               "How do we get LLMs from producing single-pass outputs to exhibiting",
               size=Pt(14), color=COLOR_BODY, bold=True)
    draw_label(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(0.6),
               "intelligent, deliberate, self-correcting behavior?",
               size=Pt(14), color=COLOR_ACCENT, bold=True)

    # Three question boxes
    questions = [
        ("Can we make the model think harder?", "→ Part 1: Internal Reasoning", COLOR_BLUE),
        ("Can we search for better solutions?", "→ Part 2: Structured Search", COLOR_GREEN),
        ("Can we ground reasoning in the real world?", "→ Part 3: Grounded Acting", COLOR_ORANGE),
    ]
    for i, (q, dest, col) in enumerate(questions):
        y = Inches(2.3) + Inches(0.7) * i
        draw_box(slide, Inches(0.8), y, Inches(4.5), Inches(0.55), q, col,
                 font_size=Pt(11), align=PP_ALIGN.LEFT)
        draw_arrow_right(slide, Inches(5.5), y + Inches(0.17), Inches(0.5), col)
        draw_label(slide, Inches(6.2), y, Inches(3.2), Inches(0.55), dest,
                   size=Pt(11), color=col, bold=True, align=PP_ALIGN.LEFT)

    # --- Slide 3: Agentic Reasoning Spectrum ---
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.placeholders[0].text = "The Agentic Reasoning Spectrum"
    for p in slide.placeholders[0].text_frame.paragraphs:
        for r in p.runs: r.font.color.rgb = COLOR_BLACK

    # Big horizontal arrow
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(0.3), Inches(2.85),
                                    Inches(9.4), Inches(0.35))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0xE8, 0xEA, 0xED)
    arrow.line.fill.background()

    stops = [
        ("IO\nPrompting", 0.2, RGBColor(0xBB, 0xDE, 0xFB)),
        ("Chain-of-\nThought", 1.7, RGBColor(0x90, 0xCA, 0xF9)),
        ("SC / LtM\nEnsemble", 3.2, RGBColor(0x64, 0xB5, 0xF6)),
        ("Tree / Graph\nSearch", 4.8, RGBColor(0x42, 0xA5, 0xF5)),
        ("ReAct\nGrounded", 6.5, RGBColor(0x21, 0x96, 0xF3)),
        ("Multi-Agent\nSystems", 8.1, RGBColor(0x19, 0x76, 0xD2)),
    ]
    for label, x, col in stops:
        draw_box(slide, Inches(x), Inches(1.6), Inches(1.3), Inches(0.9), label, col,
                 font_size=Pt(8))
        # Vertical connector line
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(x) + Inches(0.6), Inches(2.55), Inches(0.08), Inches(0.3))
        ln.fill.solid()
        ln.fill.fore_color.rgb = col
        ln.line.fill.background()

    draw_label(slide, Inches(0.3), Inches(3.4), Inches(9.4), Inches(0.4),
               "Increasing autonomy, computational cost, and capability  →",
               size=Pt(10), color=COLOR_DK2, bold=False)

    # Paper labels below
    paper_labels = [
        ("Wei '22", 0.35), ("Wang '23\nZhou '23", 1.85), ("Hao '23\nSnell '25", 3.35),
        ("Yao '23\nBesta '24\nZhou '24", 4.85), ("Yao '23\nShinn '23\nChen '23", 6.55),
        ("Kim '25", 8.3),
    ]
    for label, x in paper_labels:
        draw_label(slide, Inches(x), Inches(3.9), Inches(1.2), Inches(0.8),
                   label, size=Pt(7), color=COLOR_REF)

    # --- Slide 4: Taxonomy ---
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.placeholders[0].text = "Taxonomy: 15 Papers Across Three Pillars"
    for p in slide.placeholders[0].text_frame.paragraphs:
        for r in p.runs: r.font.color.rgb = COLOR_BLACK

    cols = [
        ("Internal Reasoning", COLOR_BLUE, COLOR_LIGHT_BLUE_BG, [
            "CoT (Wei '22)", "Self-Consistency (Wang '23)",
            "Least-to-Most (Zhou '23)", "RAP (Hao '23)",
            "Test-Time Compute (Snell '25)"
        ]),
        ("Structured Search", COLOR_GREEN, COLOR_LIGHT_GREEN_BG, [
            "Tree of Thoughts (Yao '23)", "Plan-and-Solve (Wang '23)",
            "LATS (Zhou '24)", "Graph of Thoughts (Besta '24)",
            "Planning Survey (Huang '24)", "Planning Abilities (Valmeekam '23)"
        ]),
        ("Grounded Action", COLOR_ORANGE, COLOR_LIGHT_ORANGE_BG, [
            "ReAct (Yao '23)", "Reflexion (Shinn '23)",
            "FireAct (Chen '23)", "Scaling Agents (Kim '25)"
        ]),
    ]
    for ci, (col_title, col_color, bg_color, papers) in enumerate(cols):
        x = Inches(0.3) + Inches(3.15) * ci
        # Column header
        draw_box(slide, x, Inches(1.15), Inches(2.9), Inches(0.4), col_title,
                 col_color, font_size=Pt(11))
        # Background box
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x, Inches(1.65), Inches(2.9), Inches(2.8))
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.color.rgb = col_color
        bg.line.width = Pt(1)
        try:
            bg.adjustments[0] = 0.05
        except Exception:
            pass
        # Paper list
        for pi, paper in enumerate(papers):
            draw_label(slide, x + Inches(0.1), Inches(1.75) + Inches(0.4) * pi,
                       Inches(2.7), Inches(0.35), f"• {paper}",
                       size=Pt(9), color=COLOR_BODY, align=PP_ALIGN.LEFT)

    # Arrows between columns
    for ci in range(2):
        x = Inches(3.25) + Inches(3.15) * ci
        draw_arrow_right(slide, x, Inches(2.8), Inches(0.25), COLOR_DK2)

    draw_label(slide, Inches(0.3), Inches(4.65), Inches(9.0), Inches(0.3),
               "2022  ────────────────────>  2023  ────────────────────>  2024-2025",
               size=Pt(9), color=COLOR_REF)

    # --- Slide 5: Two Dimensions ---
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.placeholders[0].text = "Two Orthogonal Dimensions of Progress"
    for p in slide.placeholders[0].text_frame.paragraphs:
        for r in p.runs: r.font.color.rgb = COLOR_BLACK

    # Y-axis label
    draw_label(slide, Inches(0.1), Inches(1.2), Inches(0.6), Inches(3.0),
               "Environmental\nGrounding\n\n↑\n\nFull Agent\n\nTool Use\n\nPure\nReasoning",
               size=Pt(8), color=COLOR_DK2, align=PP_ALIGN.CENTER)
    # X-axis label
    draw_label(slide, Inches(1.0), Inches(4.6), Inches(8.0), Inches(0.4),
               "Linear  ──────────>  Tree  ──────────>  Graph        Reasoning Topology",
               size=Pt(8), color=COLOR_DK2)

    # Place papers as dots
    paper_positions = [
        ("CoT", 1.5, 3.7, COLOR_BLUE), ("SC", 2.5, 3.5, COLOR_BLUE),
        ("LtM", 2.0, 3.2, COLOR_BLUE), ("RAP", 4.5, 3.0, COLOR_BLUE),
        ("TTC", 3.5, 2.8, COLOR_BLUE),
        ("ToT", 4.0, 3.6, COLOR_GREEN), ("PS", 2.0, 3.9, COLOR_GREEN),
        ("LATS", 5.5, 1.8, COLOR_GREEN), ("GoT", 7.0, 3.4, COLOR_GREEN),
        ("ReAct", 2.0, 1.5, COLOR_ORANGE), ("Reflexion", 2.5, 1.3, COLOR_ORANGE),
        ("FireAct", 3.0, 1.6, COLOR_ORANGE), ("Scaling", 6.5, 1.4, COLOR_ORANGE),
    ]
    for label, x, y, col in paper_positions:
        draw_box(slide, Inches(x), Inches(y), Inches(0.9), Inches(0.35), label,
                 col, font_size=Pt(7))

    # Grid lines
    for yv in [1.6, 2.5, 3.4]:
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0.8), Inches(yv), Inches(8.2), Pt(0.5))
        ln.fill.solid()
        ln.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        ln.line.fill.background()

    # --- Slide 6: Lecture Roadmap ---
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.placeholders[0].text = "Lecture Roadmap"
    for p in slide.placeholders[0].text_frame.paragraphs:
        for r in p.runs: r.font.color.rgb = COLOR_BLACK

    parts = [
        ("Part 1: Internal Reasoning", "CoT → Self-Consistency → Least-to-Most → RAP → Test-Time Compute",
         "5 papers  |  ~20 slides", COLOR_BLUE),
        ("Part 2: Structured Search & Planning", "ToT → Plan-and-Solve → LATS → GoT → Survey → Critical Analysis",
         "6 papers  |  ~24 slides", COLOR_GREEN),
        ("Part 3: Grounded Acting", "ReAct → Reflexion → FireAct → Scaling Agents",
         "4 papers  |  ~16 slides", COLOR_ORANGE),
        ("Part 4: Synthesis & Future", "Grand synthesis  |  Research trends  |  Open problems  |  Discussion",
         "~10 slides", COLOR_PURPLE),
    ]
    for i, (title, desc, meta, col) in enumerate(parts):
        y = Inches(1.2) + Inches(0.95) * i
        draw_box(slide, Inches(0.4), y, Inches(3.5), Inches(0.75), title, col,
                 font_size=Pt(11), align=PP_ALIGN.LEFT)
        draw_label(slide, Inches(4.2), y, Inches(5.0), Inches(0.45), desc,
                   size=Pt(10), color=COLOR_BODY, align=PP_ALIGN.LEFT)
        draw_label(slide, Inches(4.2), y + Inches(0.4), Inches(5.0), Inches(0.35), meta,
                   size=Pt(8), color=COLOR_REF, align=PP_ALIGN.LEFT)
        if i < 3:
            draw_arrow_right(slide, Inches(4.7), y + Inches(0.8), Inches(0.3), col)

    # --- Slide 7: Key Insight Preview ---
    add_content_slide("Key Insight Preview: Three Surprising Findings", [
        ("1. A smaller model + search can outperform a 14× larger model", True, 0),
        ("Test-Time Compute (Snell et al., 2025) shows that optimally allocating\n"
         "inference computation yields a 4× efficiency improvement over naive scaling.", False, 1),
        ("", False, 0),
        ("2. LLMs achieve only ~3% on autonomous planning tasks", True, 0),
        ("Planning Abilities (Valmeekam et al., 2023) reveals that without environment\n"
         "feedback, LLMs cannot reliably generate correct plans.", False, 1),
        ("", False, 0),
        ("3. More agents does NOT always mean better: MAS can hurt by −70%", True, 0),
        ("Scaling Agent Systems (Kim et al., 2025) shows multi-agent benefits are\n"
         "highly task-dependent — blindly adding agents can severely degrade performance.", False, 1),
    ])


# ═══════════════════════════════════════════════════════════════════
#  PART 1: INTERNAL REASONING (Slides 8-27)
# ═══════════════════════════════════════════════════════════════════

def generate_part1():
    # --- Slide 8: Section Header ---
    add_section_slide("Part 1: Internal Reasoning",
                      "Making the LLM think harder — from chains to ensembles to decomposition")

    # --- Slide 9: Section Overview Diagram ---
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.placeholders[0].text = "The Reasoning Prompting Stack"
    for p in slide.placeholders[0].text_frame.paragraphs:
        for r in p.runs: r.font.color.rgb = COLOR_BLACK

    stack = [
        ("Test-Time Compute: Optimally allocate inference budget", 0.45, Inches(1.3), RGBColor(0x0D, 0x47, 0xA1)),
        ("RAP: Plan with MCTS over reasoning steps", 0.55, Inches(1.9), RGBColor(0x15, 0x65, 0xC0)),
        ("Least-to-Most: Break down, solve piece by piece", 0.65, Inches(2.5), RGBColor(0x19, 0x76, 0xD2)),
        ("Self-Consistency: Think multiple ways, vote", 0.75, Inches(3.1), RGBColor(0x21, 0x96, 0xF3)),
        ("Chain-of-Thought: Think step by step", 0.85, Inches(3.7), RGBColor(0x42, 0xA5, 0xF5)),
    ]
    for label, w_frac, y, col in stack:
        bw = Inches(8.0) * w_frac
        x = Inches(5.0) - bw / 2
        draw_box(slide, x, y, bw, Inches(0.45), label, col, font_size=Pt(10))

    draw_label(slide, Inches(8.5), Inches(1.3), Inches(1.2), Inches(3.0),
               "↑\nIncreasing\nsophistication",
               size=Pt(8), color=COLOR_DK2, align=PP_ALIGN.CENTER)

    # ─── Paper 1: Chain-of-Thought ───
    REF_COT = "Wei et al., Chain-of-Thought Prompting Elicits Reasoning in LLMs, NeurIPS 2022"

    # Slide 10: CoT Core + Fig
    add_content_slide("Chain-of-Thought Prompting", [
        ("Core Idea", True, 0),
        ("Augment few-shot exemplars with intermediate reasoning steps", False, 1),
        ("No fine-tuning, no architectural changes — just prompt engineering", False, 1),
        ("", False, 0),
        ("Key Properties", True, 0),
        ("Decomposition: breaks multi-step problems into intermediate steps", False, 1),
        ("Interpretability: the reasoning trace is human-readable", False, 1),
        ("Generality: applicable to math, commonsense, symbolic reasoning", False, 1),
        ("Emergent ability: only works with models ≥ 100B parameters", False, 1),
    ], ref=REF_COT, fig="cot_fig3.png")

    # Slide 11: CoT Method
    add_content_slide("CoT: How It Works", [
        ("Prompt Design", True, 0),
        ("Include exemplars with <input, chain-of-thought, output> triples", False, 1),
        ("Model generates its own chain before producing the final answer", False, 1),
        ("", False, 0),
        ("Ablation Studies", True, 0),
        ("Equation-only (no natural language reasoning): degrades performance", False, 1),
        ("Reasoning-after-answer: no benefit — order matters", False, 1),
        ("Random CoT exemplars: still helps slightly vs. no CoT", False, 1),
        ("", False, 0),
        ("Why It Works", True, 0),
        ("Activates latent reasoning capabilities already present in large LMs", False, 1),
        ("Provides an \"inner monologue\" scaffold for complex multi-step tasks", False, 1),
    ], ref=REF_COT)

    # Slide 12: CoT Results + Scale Fig
    add_content_slide("CoT: Key Results", [
        ("Benchmark Performance", True, 0),
        ("GSM8K (math): 17.9% → 58.1% (PaLM 540B)", False, 1),
        ("MultiArith: 33.7% → 93.5%", False, 1),
        ("StrategyQA: 73.5% → 77.8% (+commonsense)", False, 1),
        ("", False, 0),
        ("The Emergence Phenomenon", True, 0),
        ("Models < 10B: CoT does NOT help (sometimes hurts)", False, 1),
        ("Models > 100B: CoT gives dramatic improvements", False, 1),
        ("Suggests reasoning is an emergent property of scale", False, 1),
    ], ref=REF_COT, fig="cot_fig4.png")

    # Slide 13: CoT Impact + Transition
    s = add_content_slide("CoT: Foundation for All Modern Agent Reasoning", [
        ("Impact", True, 0),
        ("Established prompt-based reasoning as a viable paradigm", False, 1),
        ("Opened the door to Tree-of-Thoughts, ReAct, and all planning methods", False, 1),
        ("Cited 10,000+ times — one of the most influential LLM papers", False, 1),
        ("", False, 0),
        ("Limitations", True, 0),
        ("Produces a single reasoning chain — what if that chain is wrong?", False, 1),
        ("No self-correction mechanism; errors propagate through the chain", False, 1),
    ], ref=REF_COT)
    add_transition_box(s, "⟶  CoT produces a single chain. But what if that chain is wrong?"
                       "  NEXT: Self-Consistency samples multiple chains and votes.", COLOR_BLUE)

    # ─── Paper 2: Self-Consistency ───
    REF_SC = "Wang et al., Self-Consistency Improves Chain of Thought Reasoning, ICLR 2023"

    # Slide 14
    add_content_slide("Self-Consistency: Sample-and-Vote Decoding", [
        ("Core Idea", True, 0),
        ("Sample multiple diverse reasoning chains via temperature sampling", False, 1),
        ("Take a majority vote over the final answers — \"self-ensemble\"", False, 1),
        ("", False, 0),
        ("Three-Step Process", True, 0),
        ("1. Prompt with CoT exemplars (same as before)", False, 1),
        ("2. Sample N diverse reasoning paths (high temperature)", False, 1),
        ("3. Marginalize over reasoning paths → majority vote on answer", False, 1),
        ("", False, 0),
        ("Key Insight", True, 0),
        ("Complex problems often have multiple valid reasoning paths", False, 1),
        ("Correct answers appear more frequently across samples", False, 1),
    ], ref=REF_SC, fig="sc_fig1.png")

    # Slide 15
    add_content_slide("Self-Consistency: Results", [
        ("Consistent Improvements Over CoT", True, 0),
        ("GSM8K: +17.9% (58.1% → 76.0%, PaLM 540B)", False, 1),
        ("ARC Challenge: +3.9%", False, 1),
        ("StrategyQA: +6.4%", False, 1),
        ("", False, 0),
        ("Scaling Behavior", True, 0),
        ("Improvement saturates around 40 sampled paths", False, 1),
        ("Works with ALL CoT methods (zero-shot, few-shot)", False, 1),
        ("Requires no additional training or fine-tuning", False, 1),
        ("", False, 0),
        ("Practical Consideration", True, 0),
        ("N× inference cost for N samples — a compute-quality tradeoff", False, 1),
    ], ref=REF_SC)

    # Slide 16
    s = add_content_slide("Self-Consistency: Impact", [
        ("Contribution", True, 0),
        ("Showed that decoding strategy matters as much as prompt design", False, 1),
        ("Established \"sample-and-aggregate\" as a general LLM improvement pattern", False, 1),
        ("Precursor to more sophisticated search methods (ToT, LATS)", False, 1),
        ("", False, 0),
        ("Limitations", True, 0),
        ("Still operates on the SAME problem — no decomposition", False, 1),
        ("Cannot handle problems that require multi-step sequential reasoning", False, 1),
    ], ref=REF_SC)
    add_transition_box(s, "⟶  SC improves decoding but still solves the SAME problem."
                       "  What if the problem itself is too complex?"
                       "  NEXT: Least-to-Most decomposes it.", COLOR_BLUE)

    # ─── Paper 3: Least-to-Most ───
    REF_LTM = "Zhou et al., Least-to-Most Prompting Enables Complex Reasoning, ICLR 2023"

    # Slide 17
    add_content_slide("Least-to-Most: Explicit Task Decomposition", [
        ("Core Idea", True, 0),
        ("Two-stage approach inspired by educational psychology:", False, 1),
        ("  Stage 1: Decompose — break complex problem into subproblems", False, 2),
        ("  Stage 2: Solve — solve each subproblem sequentially,", False, 2),
        ("  feeding prior solutions as context", False, 2),
        ("", False, 0),
        ("Key Difference from CoT", True, 0),
        ("CoT: single pass, one chain for the whole problem", False, 1),
        ("LtM: explicit decomposition → solve easy parts first → build up", False, 1),
        ("Compositional generalization to harder problems at test time", False, 1),
    ], ref=REF_LTM, fig="ltm_fig1.png")

    # Slide 18
    add_content_slide("Least-to-Most: Dramatic Results", [
        ("SCAN Benchmark (Compositional Generalization)", True, 0),
        ("CoT: 16.2% → Least-to-Most: 99.7% (!)", False, 1),
        ("Near-perfect on a task where CoT almost completely fails", False, 1),
        ("", False, 0),
        ("Last Letter Concatenation", True, 0),
        ("Trained on 2-letter problems, tested on 12-letter:", False, 1),
        ("CoT: 15.0% → Least-to-Most: 74.0%", False, 1),
        ("", False, 0),
        ("GSM8K (Math Reasoning)", True, 0),
        ("Also improves on math, though gains are more modest", False, 1),
        ("LtM shines most when compositional structure is clear", False, 1),
    ], ref=REF_LTM)

    # Slide 19: Comparison + Transition
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.placeholders[0].text = "The Reasoning Stack So Far"
    for p in slide.placeholders[0].text_frame.paragraphs:
        for r in p.runs: r.font.color.rgb = COLOR_BLACK

    methods = [
        ("Chain-of-Thought", "\"Think step by step\"", "Single path\nSingle pass", COLOR_BLUE),
        ("Self-Consistency", "\"Think multiple ways\"", "Multiple paths\nSingle pass", RGBColor(0x19, 0x76, 0xD2)),
        ("Least-to-Most", "\"Break it down\"", "Single path\nMultiple passes", RGBColor(0x0D, 0x47, 0xA1)),
    ]
    for i, (name, desc, props, col) in enumerate(methods):
        x = Inches(0.5) + Inches(3.1) * i
        draw_box(slide, x, Inches(1.3), Inches(2.8), Inches(0.45), name, col, font_size=Pt(11))
        draw_label(slide, x, Inches(1.85), Inches(2.8), Inches(0.4), desc,
                   size=Pt(10), color=COLOR_BODY, bold=True)
        draw_label(slide, x, Inches(2.3), Inches(2.8), Inches(0.7), props,
                   size=Pt(9), color=COLOR_DK2)

    add_transition_box(slide, "⟶  All three are linear or parallel. Can we SEARCH the reasoning"
                       " space systematically?  NEXT: RAP applies Monte Carlo Tree Search.",
                       COLOR_BLUE)

    # ─── Paper 4: RAP ───
    REF_RAP = "Hao et al., Reasoning with Language Model is Planning with World Model, EMNLP 2023"

    # Slide 20
    add_content_slide("RAP: LLM as World Model + Reasoning Agent", [
        ("Core Idea", True, 0),
        ("Reframe reasoning as PLANNING: LLM plays two roles simultaneously", False, 1),
        ("  World Model: predicts the next state after each reasoning step", False, 2),
        ("  Reasoning Agent: selects the best action to take", False, 2),
        ("", False, 0),
        ("Method: Monte Carlo Tree Search (MCTS)", True, 0),
        ("Selection: UCT formula balances exploration vs exploitation", False, 1),
        ("Expansion: LLM generates candidate next steps", False, 1),
        ("Simulation: LLM-as-world-model estimates future outcomes", False, 1),
        ("Backpropagation: update value estimates up the tree", False, 1),
    ], ref=REF_RAP, fig="rap_fig1.png")

    # Slide 21: MCTS figure
    add_figure_slide("RAP: Monte Carlo Tree Search in Action", "rap_fig2.png",
                     ref=REF_RAP,
                     caption="MCTS four phases: Selection → Expansion → Simulation → Backpropagation")

    # Slide 22
    add_content_slide("RAP: Results", [
        ("Blocksworld Planning", True, 0),
        ("GPT-4 (CoT): 64% → RAP (GPT-4): 95.4% (+31 points)", False, 1),
        ("LLaMA-33B (CoT): 1% → RAP: 45.3% (45× improvement!)", False, 1),
        ("", False, 0),
        ("GSM8K Math Reasoning", True, 0),
        ("CoT: 58.1% → RAP: 62.8% (with LLaMA-33B)", False, 1),
        ("More modest gains when problems are less tree-structured", False, 1),
        ("", False, 0),
        ("Key Insight", True, 0),
        ("Small model + MCTS search can match much larger models", False, 1),
        ("This foreshadows Test-Time Compute scaling results", False, 1),
    ], ref=REF_RAP)

    # Slide 23
    s = add_content_slide("RAP: Impact & Connections", [
        ("Theoretical Contribution", True, 0),
        ("First principled framework connecting LLM reasoning to classical planning", False, 1),
        ("Shows that the \"reasoning vs planning\" distinction is artificial", False, 1),
        ("", False, 0),
        ("Connections to Other Work", True, 0),
        ("ToT: uses simpler BFS/DFS instead of MCTS", False, 1),
        ("LATS: extends RAP with environment interaction + reflection", False, 1),
        ("Test-Time Compute: generalizes the compute allocation question RAP raises", False, 1),
    ], ref=REF_RAP)
    add_transition_box(s, "⟶  RAP proves search dramatically helps reasoning."
                       "  But HOW MUCH computation should we allocate?"
                       "  NEXT: Test-Time Compute provides the answer.", COLOR_BLUE)

    # ─── Paper 5: Test-Time Compute ───
    REF_TTC = "Snell et al., Scaling LLM Test-Time Compute, ICML 2025"

    # Slide 24
    add_content_slide("Scaling Test-Time Compute", [
        ("Core Finding", True, 0),
        ("Optimally allocating test-time compute yields 4× efficiency improvement", False, 1),
        ("A smaller model + optimal search can match a 14× larger model", False, 1),
        ("", False, 0),
        ("Two Knobs for Test-Time Compute", True, 0),
        ("1. Proposal distribution: revise answers iteratively (sequential refinement)", False, 1),
        ("2. Verifier: generate many candidates, pick the best (parallel search)", False, 1),
        ("", False, 0),
        ("Meta-Finding", True, 0),
        ("The optimal strategy depends on problem DIFFICULTY:", False, 1),
        ("  Easy problems → verifier-based selection is efficient", False, 2),
        ("  Hard problems → iterative revision works better", False, 2),
    ], ref=REF_TTC, fig="ttc_header.png",
       fig_pos=(Inches(5.2), Inches(1.2), Inches(4.4), Inches(1.5)))

    # Slide 25: Strategy figure
    add_figure_slide("Test-Time Compute: Strategy Comparison", "ttc_fig2.png",
                     ref=REF_TTC,
                     caption="Compute-optimal scaling: different strategies dominate at different difficulty levels")

    # Slide 26
    add_content_slide("Test-Time Compute: Key Results", [
        ("Main Results", True, 0),
        ("Compute-optimal: Llama-3.1-8B matches Llama-3.1-70B with 4× less FLOPs", False, 1),
        ("At equal compute budget: 8B + search > 70B (on many tasks)", False, 1),
        ("", False, 0),
        ("Difficulty-Dependent Strategy", True, 0),
        ("Easy: Process Reward Model (PRM) verification is optimal", False, 1),
        ("Hard: Iterative refinement outperforms parallel sampling", False, 1),
        ("Very hard: Neither helps — need fundamentally better models", False, 1),
        ("", False, 0),
        ("Connection to Real Systems", True, 0),
        ("OpenAI o1 and DeepSeek-R1 use this principle at massive scale", False, 1),
        ("\"Think longer, not bigger\" is the new scaling paradigm", False, 1),
    ], ref=REF_TTC)

    # Slide 27: Part 1 Summary + Transition
    s = add_content_slide("Part 1 Summary: The Compute-Reasoning Tradeoff", [
        ("What We Learned", True, 0),
        ("1. CoT: prompt-based reasoning is powerful but fragile (single chain)", False, 1),
        ("2. Self-Consistency: sampling multiple chains and voting improves reliability", False, 1),
        ("3. Least-to-Most: explicit decomposition enables compositional generalization", False, 1),
        ("4. RAP: tree search over reasoning achieves dramatic gains", False, 1),
        ("5. Test-Time Compute: optimal allocation of inference budget is a new scaling axis", False, 1),
        ("", False, 0),
        ("The Key Insight", True, 0),
        ("Reasoning quality scales with INFERENCE computation, not just model size", False, 1),
        ("But all methods here improve INTERNAL reasoning only...", False, 1),
    ])
    add_transition_box(s, "⟶  Part 1 improved internal reasoning."
                       "  Part 2 asks: what if we give reasoning an explicit"
                       " STRUCTURE — trees, graphs, and formal plans?", COLOR_GREEN)


# ═══════════════════════════════════════════════════════════════════
#  PART 2: STRUCTURED SEARCH & PLANNING (Slides 28-51)
# ═══════════════════════════════════════════════════════════════════

def generate_part2():
    # --- Slide 28: Section Header ---
    add_section_slide("Part 2: Structured Search & Planning",
                      "From linear chains to trees, graphs, and formal plans")

    # --- Slide 29: Topology Evolution Diagram ---
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.placeholders[0].text = "The Evolution of Reasoning Topology"
    for p in slide.placeholders[0].text_frame.paragraphs:
        for r in p.runs: r.font.color.rgb = COLOR_BLACK

    topos = [
        ("IO", "In → Out", 0.3, RGBColor(0xBD, 0xBD, 0xBD)),
        ("CoT", "In → ○ → ○ → Out", 2.0, RGBColor(0x90, 0xCA, 0xF9)),
        ("CoT-SC", "In ⇉ ○○○ → Vote", 3.9, RGBColor(0x64, 0xB5, 0xF6)),
        ("ToT", "In → Tree Search", 5.8, RGBColor(0x42, 0xA5, 0xF5)),
        ("GoT", "In → Graph Ops", 7.7, RGBColor(0x1E, 0x88, 0xE5)),
    ]
    for name, desc, x, col in topos:
        draw_box(slide, Inches(x), Inches(1.5), Inches(1.5), Inches(0.45), name, col,
                 font_size=Pt(11))
        draw_label(slide, Inches(x), Inches(2.05), Inches(1.5), Inches(0.7), desc,
                   size=Pt(8), color=COLOR_DK2)
        if x < 7.0:
            draw_arrow_right(slide, Inches(x) + Inches(1.55), Inches(1.68), Inches(0.35), COLOR_DK2)

    # LATS special: below with environment arrow
    draw_box(slide, Inches(5.8), Inches(3.2), Inches(1.5), Inches(0.45), "LATS", COLOR_GREEN,
             font_size=Pt(11))
    draw_label(slide, Inches(5.8), Inches(3.75), Inches(1.5), Inches(0.5),
               "Tree + Env\nFeedback", size=Pt(8), color=COLOR_DK2)
    draw_label(slide, Inches(0.5), Inches(3.0), Inches(5.0), Inches(0.4),
               "Adding environment interaction to search →", size=Pt(9), color=COLOR_GREEN, bold=True,
               align=PP_ALIGN.LEFT)

    # ─── Paper 6: Tree of Thoughts ───
    REF_TOT = "Yao et al., Tree of Thoughts: Deliberate Problem Solving with LLMs, NeurIPS 2023"

    # Slide 30
    add_content_slide("Tree of Thoughts (ToT)", [
        ("Core Idea", True, 0),
        ("Generalize CoT from a single chain to a TREE of reasoning paths", False, 1),
        ("LLM evaluates intermediate states → prune bad branches early", False, 1),
        ("Inspired by dual-process theory: System 1 (fast) + System 2 (deliberate)", False, 1),
        ("", False, 0),
        ("Framework: Four Design Decisions", True, 0),
        ("1. Thought decomposition: how to break problem into steps", False, 1),
        ("2. Thought generator: propose candidates (sample or propose)", False, 1),
        ("3. State evaluator: LLM judges quality of partial solutions", False, 1),
        ("4. Search algorithm: BFS or DFS over the thought tree", False, 1),
    ], ref=REF_TOT, fig="tot_fig1.png")

    # Slide 31
    add_content_slide("ToT: Design Decisions in Detail", [
        ("Thought Decomposition", True, 0),
        ("Task-specific: a \"thought\" can be a line, a sentence, or a plan step", False, 1),
        ("Must be small enough to evaluate, large enough to be meaningful", False, 1),
        ("", False, 0),
        ("State Evaluation: The Key Innovation", True, 0),
        ("LLM itself evaluates: \"Is this partial solution promising?\"", False, 1),
        ("Two modes: value (score 1-10) or vote (rank candidates)", False, 1),
        ("Enables PRUNING — skip entire subtrees that look unpromising", False, 1),
        ("", False, 0),
        ("Search Algorithms", True, 0),
        ("BFS: explore breadth-first, keep top-k at each level", False, 1),
        ("DFS: depth-first with backtracking on low-value states", False, 1),
    ], ref=REF_TOT)

    # Slide 32
    add_content_slide("ToT: Results Across Three Tasks", [
        ("Game of 24 (Mathematical Reasoning)", True, 0),
        ("IO: 7.3% → CoT: 4.0% → ToT: 74% (!!)", False, 1),
        ("CoT actually HURTS here — the problem requires exploration", False, 1),
        ("", False, 0),
        ("Creative Writing", True, 0),
        ("Generate coherent 4-paragraph passage from random sentences", False, 1),
        ("ToT (7.56) significantly outperforms IO (6.19) in human evaluation", False, 1),
        ("", False, 0),
        ("Mini Crosswords (5×5)", True, 0),
        ("IO: 15.6% words → ToT: 60% words solved", False, 1),
        ("DFS with pruning is most effective for constraint satisfaction", False, 1),
    ], ref=REF_TOT)

    # Slide 33
    s = add_content_slide("ToT: Impact", [
        ("Contribution", True, 0),
        ("Established \"LLM-as-deliberate-searcher\" paradigm", False, 1),
        ("Showed that tree search unlocks problems where linear CoT fails entirely", False, 1),
        ("Highly general: applicable whenever exploration helps", False, 1),
        ("", False, 0),
        ("Limitations", True, 0),
        ("Requires task-specific thought decomposition design", False, 1),
        ("Uses simple BFS/DFS — more sophisticated search could help", False, 1),
        ("No learning from experience — each problem solved from scratch", False, 1),
    ], ref=REF_TOT)
    add_transition_box(s, "⟶  ToT uses BFS/DFS. Can we make planning EXPLICIT"
                       " in the prompt?  NEXT: Plan-and-Solve.", COLOR_GREEN)

    # ─── Paper 7: Plan-and-Solve ───
    REF_PS = "Wang et al., Plan-and-Solve Prompting, ACL 2023"

    # Slide 34
    add_content_slide("Plan-and-Solve: Zero-Shot Planning Prompt", [
        ("Motivation", True, 0),
        ("Zero-shot CoT (\"Let's think step by step\") still makes errors:", False, 1),
        ("  Missing-step errors: skips necessary reasoning steps", False, 2),
        ("  Calculation errors: makes arithmetic mistakes", False, 2),
        ("  Semantic misunderstanding: misinterprets the question", False, 2),
        ("", False, 0),
        ("Solution: Plan-and-Solve (PS+)", True, 0),
        ("Replace \"think step by step\" with explicit planning instruction:", False, 1),
        ("\"Let's first understand the problem and devise a plan to solve it.", False, 1),
        (" Then let's carry out the plan, extract relevant variables,", False, 1),
        (" calculate intermediate results, and solve step by step.\"", False, 1),
    ], ref=REF_PS, fig="ps_fig2.png")

    # Slide 35
    add_content_slide("Plan-and-Solve: Results", [
        ("Zero-Shot Performance", True, 0),
        ("PS+ matches or exceeds few-shot CoT without any exemplars!", False, 1),
        ("GSM8K: Zero-shot PS+ (58.7%) ≈ Few-shot CoT (56.4%)", False, 1),
        ("SVAMP: PS+ (73.8%) vs CoT (69.3%)", False, 1),
        ("", False, 0),
        ("Error Reduction", True, 0),
        ("Missing-step errors: reduced by 50%+ compared to zero-shot CoT", False, 1),
        ("Calculation errors: also reduced with explicit variable tracking", False, 1),
        ("", False, 0),
        ("Significance", True, 0),
        ("Shows that plan-like structure can emerge from prompt alone", False, 1),
        ("No need for expensive few-shot exemplars", False, 1),
    ], ref=REF_PS)

    # Slide 36
    s = add_content_slide("Plan-and-Solve: Significance", [
        ("Key Takeaway", True, 0),
        ("Planning doesn't require complex tree search — even a prompt-based plan helps", False, 1),
        ("Complementary to ToT: PS plans linearly, ToT searches non-linearly", False, 1),
        ("", False, 0),
        ("Limitations", True, 0),
        ("Plan is generated in a single pass — no revision or backtracking", False, 1),
        ("Cannot handle problems where the plan itself is wrong", False, 1),
    ], ref=REF_PS)
    add_transition_box(s, "⟶  PS plans linearly. ToT searches trees."
                       "  Can we COMBINE search + planning + environmental feedback?"
                       "  NEXT: LATS unifies all three.", COLOR_GREEN)

    # ─── Paper 8: LATS ───
    REF_LATS = "Zhou et al., Language Agent Tree Search (LATS), ICML 2024"

    # Slide 37
    add_content_slide("LATS: Unifying Reasoning, Acting, and Planning", [
        ("Core Idea", True, 0),
        ("Apply Monte Carlo Tree Search to language AGENTS (not just reasoning)", False, 1),
        ("LLM serves triple duty: agent, value function, and optimizer", False, 1),
        ("", False, 0),
        ("Why LATS?", True, 0),
        ("Previous methods miss capabilities:", False, 1),
        ("  CoT/ToT: reasoning but no acting", False, 2),
        ("  ReAct: acting but no search/planning", False, 2),
        ("  RAP: search but no self-reflection", False, 2),
        ("  Reflexion: reflection but no planning", False, 2),
        ("LATS is the FIRST framework that combines ALL five capabilities", False, 1),
    ], ref=REF_LATS, fig="lats_fig1.png",
       fig_pos=(Inches(5.4), Inches(1.1), Inches(3.5), Inches(4.0)))

    # Slide 38: Comparison Table
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.placeholders[0].text = "LATS: The Complete Agent Framework"
    for p in slide.placeholders[0].text_frame.paragraphs:
        for r in p.runs: r.font.color.rgb = COLOR_BLACK

    # Build table manually with shapes
    headers = ["Method", "Reasoning", "Acting", "Planning", "Reflection", "Memory"]
    rows = [
        ("CoT", "✓", "✗", "✗", "✗", "✗"),
        ("ReAct", "✓", "✓", "✗", "✗", "✗"),
        ("ToT", "✓", "✗", "✓", "✗", "✗"),
        ("RAP", "✓", "✗", "✓", "✗", "✗"),
        ("Reflexion", "✓", "✓", "✗", "✓", "✓"),
        ("LATS", "✓", "✓", "✓", "✓", "✓"),
    ]
    col_w = Inches(1.35)
    row_h = Inches(0.38)
    x0, y0 = Inches(0.7), Inches(1.3)

    # Header row
    for ci, h in enumerate(headers):
        draw_box(slide, x0 + col_w * ci, y0, col_w - Inches(0.05), row_h,
                 h, COLOR_ACCENT, font_size=Pt(9))

    for ri, row in enumerate(rows):
        y = y0 + row_h * (ri + 1) + Inches(0.05) * (ri + 1)
        is_lats = row[0] == "LATS"
        bg = COLOR_GREEN if is_lats else COLOR_LIGHT_BG
        tc = COLOR_WHITE if is_lats else COLOR_BODY
        for ci, val in enumerate(row):
            if ci == 0:
                draw_box(slide, x0 + col_w * ci, y, col_w - Inches(0.05), row_h,
                         val, bg, text_color=tc, font_size=Pt(9))
            else:
                col = RGBColor(0x0F, 0x9D, 0x58) if val == "✓" and not is_lats else tc
                if is_lats:
                    col = COLOR_WHITE
                draw_box(slide, x0 + col_w * ci, y, col_w - Inches(0.05), row_h,
                         val, bg, text_color=col, font_size=Pt(10), bold=(val == "✓"))

    draw_label(slide, Inches(0.7), Inches(4.5), Inches(8.0), Inches(0.3),
               "LATS is the only framework that achieves all five agent capabilities.",
               size=Pt(10), color=COLOR_GREEN, bold=True, align=PP_ALIGN.LEFT)

    # Slide 39: Algorithm figure
    add_figure_slide("LATS: Six Operations in the Search Loop", "lats_algo.png",
                     ref=REF_LATS,
                     caption="Selection → Expansion → Evaluation → Simulation → Backpropagation → Reflection")

    # Slide 40
    s = add_content_slide("LATS: State-of-the-Art Results", [
        ("HotPotQA (Multi-hop QA)", True, 0),
        ("ReAct: 30% → Reflexion: 40% → LATS: 53% (GPT-3.5)", False, 1),
        ("", False, 0),
        ("HumanEval (Code Generation)", True, 0),
        ("CoT: 82.9% → LATS: 94.4% (GPT-4) — near-perfect!", False, 1),
        ("", False, 0),
        ("WebShop (Web Navigation)", True, 0),
        ("ReAct: 56.1% → LATS: 75.9% (GPT-3.5)", False, 1),
        ("", False, 0),
        ("Key Finding", True, 0),
        ("Combining search + acting + reflection synergistically outperforms each alone", False, 1),
    ], ref=REF_LATS)
    add_transition_box(s, "⟶  LATS uses trees. But what about MERGING results from different"
                       " branches?  NEXT: GoT generalizes to arbitrary graphs.", COLOR_GREEN)

    # ─── Paper 9: Graph of Thoughts ───
    REF_GOT = "Besta et al., Graph of Thoughts, AAAI 2024"

    # Slide 41
    add_content_slide("Graph of Thoughts (GoT)", [
        ("Core Idea", True, 0),
        ("Generalize from trees to arbitrary directed graphs", False, 1),
        ("Enable three novel transformations:", False, 1),
        ("  Aggregation: merge multiple partial solutions into one", False, 2),
        ("  Refinement: improve an existing thought in-place", False, 2),
        ("  Generation: branch into new reasoning paths", False, 2),
        ("", False, 0),
        ("Why Graphs > Trees?", True, 0),
        ("Trees cannot MERGE — once branches diverge, they stay separate", False, 1),
        ("Real reasoning often involves synthesizing ideas from multiple paths", False, 1),
        ("GoT = (G, T, E, R): Graph, Transformations, Evaluator, Ranker", False, 1),
    ], ref=REF_GOT, fig="got_fig1.png")

    # Slide 42
    add_content_slide("GoT: System Architecture", [
        ("Four Components", True, 0),
        ("G: the graph structure (vertices = thoughts, edges = dependencies)", False, 1),
        ("T: transformations (aggregate, refine, generate)", False, 1),
        ("E: scoring function for evaluating thought quality", False, 1),
        ("R: ranking function for selecting the best path", False, 1),
        ("", False, 0),
        ("Graph of Operations (GoO)", True, 0),
        ("User defines a static execution plan (which transformations, in what order)", False, 1),
        ("GoO is task-specific but REUSABLE across problem instances", False, 1),
        ("", False, 0),
        ("Volume of Thought", True, 0),
        ("New metric: the number of LLM calls per output token", False, 1),
        ("GoT achieves higher quality while maintaining reasonable cost", False, 1),
    ], ref=REF_GOT)

    # Slide 43
    add_content_slide("GoT: Results", [
        ("Sorting (32 elements)", True, 0),
        ("ToT: 26% improvement → GoT: 62% improvement over IO", False, 1),
        ("With 1.6× less cost than ToT", False, 1),
        ("", False, 0),
        ("Set Intersection", True, 0),
        ("GoT: 45% improvement over IO (ToT: 22%)", False, 1),
        ("Aggregation is the key: combine partial solutions from subsets", False, 1),
        ("", False, 0),
        ("Key Takeaway", True, 0),
        ("When the task benefits from MERGING intermediate results,", False, 1),
        ("graphs dramatically outperform trees", False, 1),
        ("GoT is most effective for compositional problems", False, 1),
    ], ref=REF_GOT)

    # Slide 44
    s = add_content_slide("The Full Topology Evolution", [
        ("IO → CoT → CoT-SC → ToT → GoT", True, 0),
        ("Each step adds a new structural capability:", False, 1),
        ("IO: no reasoning (baseline)", False, 2),
        ("CoT: linear chain of reasoning steps", False, 2),
        ("CoT-SC: multiple parallel chains + voting", False, 2),
        ("ToT: tree structure + pruning via LLM evaluation", False, 2),
        ("GoT: full graph + aggregation, refinement, generation", False, 2),
        ("", False, 0),
        ("Open Question", True, 0),
        ("These methods all improve the TOPOLOGY of thought.", False, 1),
        ("But can LLMs actually PLAN reliably? Two critical papers say: not really.", False, 1),
    ], ref=REF_GOT)
    add_transition_box(s, "⟶  We've seen impressive search methods."
                       "  But can LLMs actually plan?"
                       "  NEXT: Two critical analyses give a sobering answer.", COLOR_GREEN)

    # ─── Papers 10-11: Planning Survey + Planning Abilities ───
    REF_PSUR = "Huang et al., Understanding the Planning of LLM Agents: A Survey, arXiv 2024"
    REF_PA = "Valmeekam et al., On the Planning Abilities of LLMs, NeurIPS 2023"

    # Slide 45
    add_content_slide("LLM Agent Planning: A Comprehensive Taxonomy", [
        ("Five Planning Strategies (Huang et al., 2024)", True, 0),
        ("1. Task Decomposition: break goals into sub-goals (e.g., Least-to-Most)", False, 1),
        ("2. Multi-Plan Selection: generate multiple plans, pick best (e.g., ToT)", False, 1),
        ("3. External Planner-Aided: LLM generates formal spec, planner solves", False, 1),
        ("4. Reflection & Refinement: iterate on failures (e.g., Reflexion)", False, 1),
        ("5. Memory-Augmented: store past plans for retrieval (e.g., LATS)", False, 1),
    ], ref=REF_PSUR, fig="psur_fig1.png")

    # Slide 46
    add_content_slide("Mapping Our Lecture Papers to the Taxonomy", [
        ("Task Decomposition", True, 0),
        ("Least-to-Most, Plan-and-Solve", False, 1),
        ("", False, 0),
        ("Multi-Plan Selection", True, 0),
        ("Self-Consistency, ToT, RAP, GoT", False, 1),
        ("", False, 0),
        ("External Planner-Aided", True, 0),
        ("Not covered in depth — involves PDDL solvers, external tools", False, 1),
        ("", False, 0),
        ("Reflection & Refinement", True, 0),
        ("Reflexion, LATS", False, 1),
        ("", False, 0),
        ("Memory-Augmented", True, 0),
        ("LATS (stores experience for future search), Reflexion (episodic memory)", False, 1),
    ], ref=REF_PSUR)

    # Slide 47
    add_content_slide("Planning Survey: Benchmark Landscape", [
        ("Interactive Benchmarks", True, 0),
        ("ALFWorld: household tasks in text environment", False, 1),
        ("WebShop: online shopping navigation", False, 1),
        ("ScienceWorld: scientific experiments simulation", False, 1),
        ("", False, 0),
        ("QA / Reasoning Benchmarks", True, 0),
        ("HotPotQA: multi-hop question answering", False, 1),
        ("FEVER: fact verification", False, 1),
        ("GSM8K / MATH: mathematical reasoning", False, 1),
        ("", False, 0),
        ("Code Benchmarks", True, 0),
        ("HumanEval: code generation", False, 1),
        ("MBPP: basic Python programming", False, 1),
    ], ref=REF_PSUR)

    # Slide 48: Critical finding
    add_content_slide("Can LLMs Really Plan? A Critical Investigation", [
        ("Valmeekam et al. (NeurIPS 2023): A Sobering Finding", True, 0),
        ("", False, 0),
        ("Autonomous Planning Success Rate: ~3%", True, 0),
        ("GPT-4 achieves only 3% on standard planning benchmarks (Blocksworld)", False, 1),
        ("When LLMs must generate complete plans from scratch, they fail badly", False, 1),
        ("", False, 0),
        ("BUT: Plans Are Repairable", True, 0),
        ("LLMs can identify errors in given plans with ~40% accuracy", False, 1),
        ("Iterative correction: plans improve with each round of feedback", False, 1),
        ("This explains why ReAct-style approaches work — environment corrects!", False, 1),
    ], ref=REF_PA, fig="pa_fig1.png",
       fig_pos=(Inches(5.4), Inches(1.5), Inches(4.2), Inches(2.5)))

    # Slide 49
    s = add_content_slide("Planning Abilities: Implications for Agent Architecture", [
        ("The Key Tension", True, 0),
        ("ReAct says: LLMs plan well (when grounded in environment feedback)", False, 1),
        ("Valmeekam says: LLMs cannot plan (autonomous, from scratch)", False, 1),
        ("", False, 0),
        ("Resolution", True, 0),
        ("LLMs are not autonomous planners — they are good COLLABORATORS", False, 1),
        ("Environment feedback, tool use, and verification correct their errors", False, 1),
        ("This is exactly why grounded acting matters so much", False, 1),
        ("", False, 0),
        ("Design Principle", True, 0),
        ("Never let LLMs plan in isolation — always provide feedback loops", False, 1),
        ("The best architectures combine LLM generation with external verification", False, 1),
    ], ref=REF_PA)
    add_transition_box(s, "⟶  Part 2 reveals a tension: LLMs struggle with pure planning"
                       " but excel when grounded."
                       "  Part 3 explores how to ground reasoning in the real world.", COLOR_ORANGE)

    # --- Slides 50-51: Part 2 Synthesis ---
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.placeholders[0].text = "Part 2 Synthesis: The Search-Planning Landscape"
    for p in slide.placeholders[0].text_frame.paragraphs:
        for r in p.runs: r.font.color.rgb = COLOR_BLACK

    # 2D scatter using shapes
    draw_label(slide, Inches(0.1), Inches(1.2), Inches(0.7), Inches(3.3),
               "Env.\nInteraction\n\n↑\n\nFull\n\nPartial\n\nNone",
               size=Pt(7), color=COLOR_DK2)
    draw_label(slide, Inches(1.0), Inches(4.5), Inches(8.0), Inches(0.4),
               "Linear  ──────────>  Tree  ──────────>  Graph        Search Complexity",
               size=Pt(8), color=COLOR_DK2)

    papers_2d = [
        ("PS", 1.5, 3.6, COLOR_GREEN), ("ToT", 4.0, 3.4, COLOR_GREEN),
        ("GoT", 7.5, 3.2, COLOR_GREEN), ("RAP", 4.5, 2.8, COLOR_BLUE),
        ("LATS", 5.0, 1.6, COLOR_GREEN),
    ]
    for name, x, y, col in papers_2d:
        draw_box(slide, Inches(x), Inches(y), Inches(1.0), Inches(0.35), name, col, font_size=Pt(9))

    # Warning zone
    warn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(1.2), Inches(1.3), Inches(7.0), Inches(0.5))
    warn.fill.solid()
    warn.fill.fore_color.rgb = RGBColor(0xFC, 0xE4, 0xEC)
    warn.line.color.rgb = COLOR_RED
    warn.line.width = Pt(1)
    try:
        warn.adjustments[0] = 0.2
    except Exception:
        pass
    tf = warn.text_frame
    _add_run(tf.paragraphs[0], "⚠ Warning (Valmeekam): LLMs alone achieve ~3% on autonomous planning",
             size=Pt(9), color=COLOR_RED, bold=True, font_name=FONT_BODY)

    # Slide 51: Transition
    s = add_content_slide("From Thinking to Acting", [
        ("Part 2 Key Findings", True, 0),
        ("1. Tree search (ToT) unlocks problems where linear CoT fails (4% → 74%)", False, 1),
        ("2. Even a prompt-level plan (PS) helps match few-shot performance", False, 1),
        ("3. LATS combines ALL agent capabilities (reasoning + acting + planning + reflection)", False, 1),
        ("4. GoT enables merging partial solutions — graphs > trees for some tasks", False, 1),
        ("5. But LLMs alone cannot plan reliably (~3% autonomous success)", False, 1),
        ("", False, 0),
        ("The Critical Implication", True, 0),
        ("Reasoning alone is not enough. Agents must ACT in the world,", False, 1),
        ("receive feedback, and LEARN from their mistakes.", False, 1),
    ])
    add_transition_box(s, "⟶  Part 3: How do we connect reasoning to the real world?"
                       "  ReAct, Reflexion, FireAct, and Scaling Agent Systems.", COLOR_ORANGE)


# ═══════════════════════════════════════════════════════════════════
#  PART 3: GROUNDED ACTING (Slides 52-68)
# ═══════════════════════════════════════════════════════════════════

def generate_part3():
    # --- Slide 52: Section Header ---
    add_section_slide("Part 3: Grounded Acting & ReAct",
                      "Connecting reasoning to the real world")

    # --- Slide 53: Agent Evolution Diagram ---
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.placeholders[0].text = "The Agent Evolution"
    for p in slide.placeholders[0].text_frame.paragraphs:
        for r in p.runs: r.font.color.rgb = COLOR_BLACK

    evos = [
        ("ReAct", "The Basic Agent", "Thought → Action\n→ Observation loop", Inches(0.5), COLOR_ORANGE),
        ("Reflexion", "The Learning Agent", "+ Self-critique\n+ Episodic memory", Inches(2.7), RGBColor(0xE6, 0x5C, 0x00)),
        ("FireAct", "The Efficient Agent", "+ Fine-tuning on\nagent trajectories", Inches(4.9), RGBColor(0xBF, 0x36, 0x0C)),
        ("Scaling", "The System", "When to use\nmultiple agents?", Inches(7.1), COLOR_RED),
    ]
    for name, subtitle, desc, x, col in evos:
        draw_box(slide, x, Inches(1.5), Inches(1.8), Inches(0.5), name, col, font_size=Pt(12))
        draw_label(slide, x, Inches(2.1), Inches(1.8), Inches(0.35), subtitle,
                   size=Pt(9), color=col, bold=True)
        draw_label(slide, x, Inches(2.5), Inches(1.8), Inches(0.7), desc,
                   size=Pt(8), color=COLOR_DK2)
        if x < Inches(6):
            draw_arrow_right(slide, x + Inches(1.85), Inches(1.7), Inches(0.75), COLOR_DK2)

    draw_label(slide, Inches(0.5), Inches(3.5), Inches(8.5), Inches(0.5),
               "Each step adds a new capability: ReAct grounds reasoning,"
               " Reflexion adds learning, FireAct adds training, Scaling adds system design.",
               size=Pt(10), color=COLOR_BODY, align=PP_ALIGN.LEFT)

    # ─── Paper 12: ReAct ───
    REF_REACT = "Yao et al., ReAct: Synergizing Reasoning and Acting in LLMs, ICLR 2023"

    # Slide 54
    add_content_slide("ReAct: Thought-Action-Observation Loop", [
        ("Core Idea", True, 0),
        ("Interleave reasoning TRACES with concrete ACTIONS in the environment", False, 1),
        ("The agent thinks (Thought), acts (Action), observes (Observation), and repeats", False, 1),
        ("", False, 0),
        ("Why Interleave?", True, 0),
        ("Reasoning-only (CoT): 14% hallucinated facts, 56% hallucinates in failure cases", False, 1),
        ("ReAct: 0% hallucinated facts (environment keeps the agent grounded)", False, 1),
        ("Acting-only: makes repetitive, goalless actions without reasoning", False, 1),
        ("", False, 0),
        ("Key Innovation", True, 0),
        ("\"Thought\" is a FREE-FORM reasoning trace (not a fixed tool call)", False, 1),
        ("Enables the agent to plan, track progress, handle exceptions", False, 1),
    ], ref=REF_REACT, fig="react_fig1.png")

    # Slide 55
    add_content_slide("ReAct: Method", [
        ("Augmented Action Space", True, 0),
        ("â = Action ∪ Thought (reasoning is treated as a special action)", False, 1),
        ("Thoughts don't change the environment but update the agent's reasoning", False, 1),
        ("", False, 0),
        ("The ReAct Loop", True, 0),
        ("1. Thought: \"I need to search for X to find Y\"", False, 1),
        ("2. Action: search[X]", False, 1),
        ("3. Observation: [search results from environment]", False, 1),
        ("4. Thought: \"The results show Z, so I should now...\"", False, 1),
        ("5. Repeat until task is solved or max steps reached", False, 1),
        ("", False, 0),
        ("Composability", True, 0),
        ("ReAct + Self-Consistency: sample multiple ReAct trajectories, vote → best results", False, 1),
    ], ref=REF_REACT)

    # Slide 56
    add_content_slide("ReAct: Results Across Four Benchmarks", [
        ("Knowledge-Intensive Reasoning", True, 0),
        ("HotPotQA: Act-only 25.7% → ReAct 27.4% → ReAct+SC 35.1%", False, 1),
        ("FEVER: Act-only 55.7% → ReAct 60.9% → ReAct+SC 64.6%", False, 1),
        ("", False, 0),
        ("Interactive Decision Making", True, 0),
        ("ALFWorld: BUTLER (10%) → ReAct (71%) — 7× improvement!", False, 1),
        ("WebShop: IL+RL (29.1%) → ReAct (40.0%)", False, 1),
        ("", False, 0),
        ("Error Analysis: Why ReAct Wins", True, 0),
        ("CoT failure: 56% involve hallucinated facts", False, 1),
        ("ReAct failure: 0% hallucination — errors are about search strategy, not facts", False, 1),
    ], ref=REF_REACT)

    # Slide 57
    s = add_content_slide("ReAct: THE Foundational Agent Paradigm", [
        ("Impact", True, 0),
        ("The most widely adopted agent architecture in the LLM era", False, 1),
        ("Used by: LangChain, AutoGPT, BabyAGI, ChatGPT Plugins, Claude Tools", False, 1),
        ("Cited 5000+ times — defined \"what an LLM agent looks like\"", False, 1),
        ("", False, 0),
        ("Limitations", True, 0),
        ("No learning from past mistakes — each episode starts fresh", False, 1),
        ("Performance degrades over long horizons (error accumulation)", False, 1),
        ("Depends heavily on prompt engineering for the Thought traces", False, 1),
    ], ref=REF_REACT)
    add_transition_box(s, "⟶  ReAct doesn't learn from failures."
                       "  What if the agent could REFLECT on its mistakes?"
                       "  NEXT: Reflexion.", COLOR_ORANGE)

    # ─── Paper 13: Reflexion ───
    REF_REFL = "Shinn et al., Reflexion: Language Agents with Verbal Reinforcement Learning, NeurIPS 2023"

    # Slide 58
    add_content_slide("Reflexion: Learning from Verbal Self-Critique", [
        ("Core Idea", True, 0),
        ("\"Verbal reinforcement learning\": agent reflects on failures in natural language", False, 1),
        ("Store reflections in memory → use them to improve future attempts", False, 1),
        ("No weight updates needed — learning happens through prompting", False, 1),
        ("", False, 0),
        ("Three Components", True, 0),
        ("1. Actor: generates trajectories using ReAct-style reasoning", False, 1),
        ("2. Evaluator: assesses whether the trajectory succeeded", False, 1),
        ("3. Self-Reflection: generates verbal critique of failed attempts", False, 1),
    ], ref=REF_REFL, fig="refl_fig2.png")

    # Slide 59
    add_content_slide("Reflexion: Architecture & Loop", [
        ("The Reflexion Loop", True, 0),
        ("1. Actor attempts the task (using ReAct)", False, 1),
        ("2. Evaluator checks: success or failure?", False, 1),
        ("3. If failed: Self-Reflection generates a verbal critique", False, 1),
        ("   \"I searched for X but should have searched for Y because...\"", False, 2),
        ("4. Critique is stored in episodic memory", False, 1),
        ("5. Next attempt: critique is included in the prompt as guidance", False, 1),
        ("", False, 0),
        ("Key Properties", True, 0),
        ("No gradient updates — purely prompt-based learning", False, 1),
        ("Memory grows with each trial — agent gets smarter over time", False, 1),
        ("Self-critique is more informative than scalar reward signals", False, 1),
    ], ref=REF_REFL)

    # Slide 60
    add_content_slide("Reflexion: Results", [
        ("ALFWorld (Household Tasks)", True, 0),
        ("ReAct: 75% → Reflexion: 97% (near-perfect with 2-3 trials)", False, 1),
        ("", False, 0),
        ("HotPotQA (Multi-hop QA)", True, 0),
        ("ReAct: 30% → Reflexion: 40% → Reflexion+EPM: 51%", False, 1),
        ("EPM = Episodic Planning Memory for additional context", False, 1),
        ("", False, 0),
        ("HumanEval (Code Generation)", True, 0),
        ("GPT-4 (pass@1): 80.1% → Reflexion: 91.0%", False, 1),
        ("Achieves SotA on code generation at the time of publication", False, 1),
        ("", False, 0),
        ("Convergence", True, 0),
        ("Most tasks solved within 3-5 reflection rounds", False, 1),
    ], ref=REF_REFL)

    # Slide 61
    s = add_content_slide("Reflexion: Impact on Agentic AI", [
        ("Contribution", True, 0),
        ("Pioneered \"verbal learning\" — agents improve without retraining", False, 1),
        ("Showed that natural language memory is a viable learning mechanism", False, 1),
        ("Connects to human metacognition — learning from experience", False, 1),
        ("", False, 0),
        ("Limitations", True, 0),
        ("Memory is session-specific — doesn't persist across tasks", False, 1),
        ("Self-critique quality depends on the base model's capability", False, 1),
        ("Cannot learn fundamentally new skills — only refine existing ones", False, 1),
    ], ref=REF_REFL)
    add_transition_box(s, "⟶  Reflexion improves through prompting."
                       "  Can we go further and TRAIN better agents?"
                       "  NEXT: FireAct fine-tunes on agent trajectories.", COLOR_ORANGE)

    # ─── Paper 14: FireAct ───
    REF_FA = "Chen et al., FireAct: Toward Language Agent Fine-tuning, arXiv 2023"

    # Slide 62
    add_content_slide("FireAct: Fine-tuning Language Agents", [
        ("Core Idea", True, 0),
        ("Fine-tune smaller LLMs on agent TRAJECTORIES, not just QA pairs", False, 1),
        ("Multi-method training: mix ReAct + CoT + Reflexion trajectories", False, 1),
        ("", False, 0),
        ("Training Pipeline", True, 0),
        ("1. Generate successful trajectories using GPT-4 with multiple methods", False, 1),
        ("2. Filter: keep only trajectories that lead to correct answers", False, 1),
        ("3. Fine-tune smaller model (Llama-2-13B) on this diverse mixture", False, 1),
        ("", False, 0),
        ("Key Innovation", True, 0),
        ("Method DIVERSITY in training data matters more than volume", False, 1),
        ("Agent trained on mixed trajectories outperforms any single method", False, 1),
    ], ref=REF_FA, fig="fa_fig.png",
       fig_pos=(Inches(5.4), Inches(1.5), Inches(4.2), Inches(2.3)))

    # Slide 63
    add_content_slide("FireAct: Results & Efficiency", [
        ("Main Results (HotPotQA)", True, 0),
        ("GPT-3.5 + ReAct prompting: 27.4%", False, 1),
        ("Llama-2-13B + FireAct fine-tuning: 31.8% (outperforms GPT-3.5!)", False, 1),
        ("GPT-4 + ReAct prompting: 35.1% (FireAct approaches this)", False, 1),
        ("", False, 0),
        ("Efficiency Gains", True, 0),
        ("70% reduction in inference time vs. prompted agents", False, 1),
        ("Smaller model = cheaper, faster, more deployable", False, 1),
        ("", False, 0),
        ("Key Finding: Trajectory Diversity", True, 0),
        ("Training on MIXED methods (ReAct+CoT+Reflexion) > any single method", False, 1),
        ("Diverse reasoning styles produce a more robust agent", False, 1),
    ], ref=REF_FA)

    # Slide 64
    s = add_content_slide("FireAct: Bridging Prompting and Training", [
        ("Significance", True, 0),
        ("First systematic study of fine-tuning for language AGENTS (not just LLMs)", False, 1),
        ("Demonstrates the prompting → training pipeline for agent improvement", False, 1),
        ("Shows small models can match large prompted models as agents", False, 1),
        ("", False, 0),
        ("Key Takeaway", True, 0),
        ("Trajectory diversity > trajectory volume", False, 1),
        ("The future of agent improvement: generate diverse trajectories, then distill", False, 1),
    ], ref=REF_FA)
    add_transition_box(s, "⟶  FireAct improves individual agents."
                       "  But should we use ONE strong agent or MANY?"
                       "  NEXT: Scaling Agent Systems.", COLOR_ORANGE)

    # ─── Paper 15: Scaling Agent Systems ───
    REF_SA = "Kim et al., Towards a Science of Scaling Agent Systems, arXiv 2025 (Google Research)"

    # Slide 65
    add_content_slide("Scaling Agent Systems: More Agents ≠ Better", [
        ("Core Finding", True, 0),
        ("Multi-agent systems (MAS) can improve performance by +81%...", False, 1),
        ("...or HURT performance by −70%! Benefits are task-contingent.", False, 1),
        ("", False, 0),
        ("Study Design", True, 0),
        ("Five architectures: SAS, MAS-Independent, MAS-Decentralized,", False, 1),
        ("  MAS-Centralized, MAS-Hybrid", False, 2),
        ("Three model families: Gemini, GPT, Claude", False, 1),
        ("Five benchmarks: GSM8K, MATH, GPQA, HumanEval, MBPP", False, 1),
    ], ref=REF_SA, fig="sa_fig1.png")

    # Slide 66: Architecture figure
    add_figure_slide("Five Canonical Agent Architectures", "sa_fig3.png",
                     ref=REF_SA,
                     caption="SAS = Single Agent | MAS-I = Independent | MAS-D = Decentralized | MAS-C = Centralized | MAS-H = Hybrid")

    # Slide 67
    add_content_slide("Three Fundamental Scaling Principles", [
        ("Principle 1: Tool-Coordination Trade-off", True, 0),
        ("Tasks needing tools (code, math): MAS helps by parallelizing tool use", False, 1),
        ("Tasks needing knowledge: MAS hurts by fragmenting reasoning context", False, 1),
        ("", False, 0),
        ("Principle 2: Capability Saturation", True, 0),
        ("Strong models (GPT-4, Gemini-Pro) gain LESS from multi-agent scaling", False, 1),
        ("Weaker models benefit more — MAS compensates for individual limitations", False, 1),
        ("", False, 0),
        ("Principle 3: Error Amplification via Topology", True, 0),
        ("Centralized: single point of failure at the coordinator", False, 1),
        ("Decentralized: error propagation along communication chains", False, 1),
        ("Hybrid: balances both failure modes", False, 1),
    ], ref=REF_SA)

    # Slide 68
    add_content_slide("When to Use Multi-Agent Systems", [
        ("Decision Framework (87% prediction accuracy)", True, 0),
        ("", False, 0),
        ("USE multi-agent when:", True, 0),
        ("Task requires heavy tool use (code execution, search)", False, 1),
        ("Individual model is weak and benefits from collaboration", False, 1),
        ("Task is naturally parallelizable into independent subtasks", False, 1),
        ("", False, 0),
        ("AVOID multi-agent when:", True, 0),
        ("Task requires deep, coherent reasoning (knowledge-intensive QA)", False, 1),
        ("Individual model is already strong (GPT-4, Gemini-Pro)", False, 1),
        ("Communication overhead > collaboration benefit", False, 1),
        ("", False, 0),
        ("Bottom Line", True, 0),
        ("\"More agents\" is not a universal solution — architecture must match the task", False, 1),
    ], ref=REF_SA)


# ═══════════════════════════════════════════════════════════════════
#  PART 4: SYNTHESIS & RESEARCH DIRECTIONS (Slides 69-78)
# ═══════════════════════════════════════════════════════════════════

def generate_part4():
    # --- Slide 69: Grand Synthesis ---
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.placeholders[0].text = "Grand Synthesis: How It All Connects"
    for p in slide.placeholders[0].text_frame.paragraphs:
        for r in p.runs: r.font.color.rgb = COLOR_BLACK

    # Network graph of papers
    nodes = [
        ("CoT\n(Wei '22)", 0.5, 3.5, COLOR_BLUE),
        ("SC\n(Wang '23)", 0.5, 1.8, COLOR_BLUE),
        ("LtM\n(Zhou '23)", 2.3, 3.5, COLOR_BLUE),
        ("RAP\n(Hao '23)", 2.3, 1.8, COLOR_BLUE),
        ("TTC\n(Snell '25)", 2.3, 0.2, COLOR_BLUE),
        ("ToT\n(Yao '23)", 4.1, 3.2, COLOR_GREEN),
        ("PS\n(Wang '23)", 4.1, 1.5, COLOR_GREEN),
        ("LATS\n(Zhou '24)", 5.9, 2.3, COLOR_GREEN),
        ("GoT\n(Besta '24)", 5.9, 0.5, COLOR_GREEN),
        ("ReAct\n(Yao '23)", 4.1, 4.8, COLOR_ORANGE),
        ("Reflexion\n(Shinn '23)", 5.9, 4.8, COLOR_ORANGE),
        ("FireAct\n(Chen '23)", 7.7, 4.8, COLOR_ORANGE),
        ("Scaling\n(Kim '25)", 7.7, 2.3, COLOR_RED),
    ]
    for label, x, y_off, col in nodes:
        draw_box(slide, Inches(x), Inches(1.2 + y_off * 0.65), Inches(1.4), Inches(0.5),
                 label, col, font_size=Pt(7))

    # Simplified connections via arrows (key influences only)
    draw_label(slide, Inches(0.3), Inches(4.9), Inches(9.0), Inches(0.3),
               "Arrows show key influences: CoT → ToT/ReAct/SC  |  ToT → LATS/GoT  |  "
               "ReAct → Reflexion → FireAct  |  TTC meta-analyzes compute allocation",
               size=Pt(7), color=COLOR_REF, align=PP_ALIGN.LEFT)

    # --- Slide 70: Key Tensions ---
    add_content_slide("Three Key Tensions in the Field", [
        ("Tension 1: Can LLMs Plan?", True, 0),
        ("Valmeekam (NeurIPS '23): NO — only ~3% success on autonomous planning", False, 1),
        ("ReAct/LATS: YES — when grounded in environment feedback", False, 1),
        ("Resolution: LLMs are collaborators, not autonomous planners", False, 2),
        ("", False, 0),
        ("Tension 2: Scale Model or Scale Inference?", True, 0),
        ("Traditional scaling: bigger model = better performance", False, 1),
        ("Test-Time Compute: smaller model + search = matching larger models", False, 1),
        ("Resolution: optimal strategy depends on problem difficulty", False, 2),
        ("", False, 0),
        ("Tension 3: One Agent or Many?", True, 0),
        ("Intuition: more agents should help (division of labor)", False, 1),
        ("Reality: MAS can hurt by −70% (Kim et al., 2025)", False, 1),
        ("Resolution: architecture must match task structure", False, 2),
    ])

    # --- Slide 71: Research Direction 1 ---
    add_content_slide("Trend 1: Test-Time Compute as the New Scaling Axis", [
        ("The Shift", True, 0),
        ("From: \"Train a bigger model\" (pre-training scaling)", False, 1),
        ("To: \"Think longer at inference\" (test-time scaling)", False, 1),
        ("", False, 0),
        ("Evidence", True, 0),
        ("OpenAI o1/o3: chain-of-thought reasoning at massive test-time compute", False, 1),
        ("DeepSeek-R1: open-weight model with inference-time reasoning scaling", False, 1),
        ("Snell et al.: formal framework for compute-optimal allocation", False, 1),
        ("", False, 0),
        ("Open Questions", True, 0),
        ("How to build adaptive routing: fast (System 1) vs slow (System 2) thinking?", False, 1),
        ("What is the theoretical limit of test-time compute scaling?", False, 1),
        ("Can we train verifiers that are more reliable than the generators?", False, 1),
    ])

    # --- Slide 72: Research Direction 2 ---
    add_content_slide("Trend 2: From Prompting to Training for Agents", [
        ("The Evolution", True, 0),
        ("2022-23: Prompt engineering (CoT, ReAct, Reflexion)", False, 1),
        ("2024: Fine-tuning on trajectories (FireAct)", False, 1),
        ("2025+: RLHF/RLAIF for agent behavior optimization", False, 1),
        ("", False, 0),
        ("Emerging Approaches", True, 0),
        ("Agent distillation: GPT-4 generates trajectories → train smaller agents", False, 1),
        ("Online learning: agents improve during deployment (Reflexion-style)", False, 1),
        ("Multi-task agent training: one model, many tools and environments", False, 1),
        ("", False, 0),
        ("Challenge", True, 0),
        ("Training data = (state, action, reward) trajectories are expensive", False, 1),
        ("How to learn from imperfect/noisy/contradictory demonstrations?", False, 1),
    ])

    # --- Slide 73: Research Direction 3 ---
    add_content_slide("Trend 3: Hybrid Architectures", [
        ("The Problem with Pure LLM Planning", True, 0),
        ("Valmeekam: ~3% success on autonomous planning", False, 1),
        ("LLMs are stochastic — not suitable for guaranteed-correct plans", False, 1),
        ("", False, 0),
        ("Neuro-Symbolic Hybrid Approach", True, 0),
        ("LLM generates formal specification (e.g., PDDL domain/problem)", False, 1),
        ("Classical planner produces guaranteed-correct plan", False, 1),
        ("LLM interprets and executes the plan in natural language", False, 1),
        ("", False, 0),
        ("Other Hybrid Patterns", True, 0),
        ("LLM + Code Interpreter: generate code → execute → verify results", False, 1),
        ("LLM + Retrieval: ground reasoning in retrieved facts (RAG)", False, 1),
        ("LLM + Formal Verifier: generate proofs → check with theorem prover", False, 1),
    ])

    # --- Slide 74: Research Direction 4 ---
    add_content_slide("Trend 4: Self-Improving Agents", [
        ("From Reflexion to Continual Learning", True, 0),
        ("Reflexion: verbal memory for within-task improvement", False, 1),
        ("CLIN (NeurIPS 2024): cross-task knowledge transfer via memory", False, 1),
        ("Future: agents that accumulate skills over deployment lifetime", False, 1),
        ("", False, 0),
        ("Key Challenges", True, 0),
        ("Long-horizon memory: how to store and retrieve relevant experience?", False, 1),
        ("Catastrophic forgetting: new learning shouldn't erase old skills", False, 1),
        ("Safety: self-improving agents must not learn harmful behaviors", False, 1),
        ("", False, 0),
        ("Vision", True, 0),
        ("Agents that get better with every interaction — like human experts", False, 1),
        ("\"The best agent is not the biggest model, but the one with the most experience\"", False, 1),
    ])

    # --- Slide 75: Open Problems ---
    add_content_slide("Open Problems in Agentic AI", [
        ("1. Reliability & Safety", True, 0),
        ("How to ensure agents don't take harmful actions in the real world?", False, 1),
        ("Hallucination → incorrect tool calls → cascading failures", False, 1),
        ("", False, 0),
        ("2. Evaluation", True, 0),
        ("Current benchmarks test narrow tasks — how to evaluate general agency?", False, 1),
        ("Need: long-horizon, multi-domain, open-ended evaluation suites", False, 1),
        ("", False, 0),
        ("3. Cost-Performance Tradeoffs", True, 0),
        ("Test-time compute: diminishing returns on very hard problems", False, 1),
        ("Multi-agent: overhead can exceed benefits", False, 1),
        ("", False, 0),
        ("4. Grounding & Embodiment", True, 0),
        ("Moving from text environments to physical world interaction", False, 1),
        ("Sim-to-real transfer for embodied agents", False, 1),
    ])

    # --- Slide 76: Discussion Questions ---
    add_content_slide("Discussion Questions", [
        ("1. The Planning Paradox", True, 0),
        ("If LLMs only achieve 3% on autonomous planning, should we abandon"
         " LLM-based planning entirely? Or is environment-grounded planning sufficient?", False, 1),
        ("", False, 0),
        ("2. The Scaling Question", True, 0),
        ("Given Test-Time Compute results, is training ever-larger models still"
         " the right approach? Or should we invest in better inference algorithms?", False, 1),
        ("", False, 0),
        ("3. Agent Architecture Design", True, 0),
        ("For a real-world application (e.g., autonomous coding assistant), which"
         " combination of methods from this lecture would you use and why?", False, 1),
        ("", False, 0),
        ("4. The Multi-Agent Debate", True, 0),
        ("Kim et al. show MAS can hurt. Under what conditions should a practitioner"
         " choose multi-agent over single-agent, and how would you decide?", False, 1),
        ("", False, 0),
        ("5. Self-Improvement Safety", True, 0),
        ("How should we constrain self-improving agents (Reflexion-style) to prevent"
         " them from learning unsafe behaviors through trial and error?", False, 1),
    ], bullet_size=Pt(10))

    # --- Slide 77: Reading List ---
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.placeholders[0].text = "Reading List: W2 Papers"
    for p in slide.placeholders[0].text_frame.paragraphs:
        for r in p.runs: r.font.color.rgb = COLOR_BLACK

    papers_text = [
        ("Part 1: Internal Reasoning", True, 0),
        ("[1] Wei et al. (2022). Chain-of-Thought Prompting. NeurIPS", False, 1),
        ("[2] Wang et al. (2023). Self-Consistency. ICLR", False, 1),
        ("[3] Zhou et al. (2023). Least-to-Most Prompting. ICLR", False, 1),
        ("[4] Hao et al. (2023). RAP: Reasoning with Language Model is Planning. EMNLP", False, 1),
        ("[5] Snell et al. (2025). Scaling LLM Test-Time Compute. ICML", False, 1),
        ("Part 2: Structured Search & Planning", True, 0),
        ("[6] Yao et al. (2023). Tree of Thoughts. NeurIPS", False, 1),
        ("[7] Wang et al. (2023). Plan-and-Solve Prompting. ACL", False, 1),
        ("[8] Zhou et al. (2024). LATS. ICML", False, 1),
        ("[9] Besta et al. (2024). Graph of Thoughts. AAAI", False, 1),
        ("[10] Huang et al. (2024). Planning Survey. arXiv", False, 1),
        ("[11] Valmeekam et al. (2023). Planning Abilities. NeurIPS", False, 1),
        ("Part 3: Grounded Acting", True, 0),
        ("[12] Yao et al. (2023). ReAct. ICLR", False, 1),
        ("[13] Shinn et al. (2023). Reflexion. NeurIPS", False, 1),
        ("[14] Chen et al. (2023). FireAct. arXiv", False, 1),
        ("[15] Kim et al. (2025). Scaling Agent Systems. arXiv / Google Research", False, 1),
    ]
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(1.15), Inches(9.0), Inches(3.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, indent) in enumerate(papers_text):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _add_run(para, text, size=Pt(8), bold=bold,
                 color=COLOR_ACCENT if bold else COLOR_BODY)
        para.level = indent
        para.space_after = Pt(2)

    # --- Slide 78: Thank You ---
    add_title_slide("Thank You", "AI 89900 — Agentic AI  |  KAIST Spring 2026")


# ═══════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    print("Generating W2 slides...")
    print("  Part 0: Opening (7 slides)")
    generate_part0()
    print("  Part 1: Internal Reasoning (20 slides)")
    generate_part1()
    print("  Part 2: Structured Search & Planning (24 slides)")
    generate_part2()
    print("  Part 3: Grounded Acting (17 slides)")
    generate_part3()
    print("  Part 4: Synthesis & Research Directions (10 slides)")
    generate_part4()

    total = len(prs.slides)
    print(f"\n  Total: {total} slides")

    out_path = "/Users/sungju/Agentic AI/slides/W2_Agentic_Architecture_Patterns.pptx"
    prs.save(out_path)
    print(f"  Saved: {out_path}")
    print("Done!")
