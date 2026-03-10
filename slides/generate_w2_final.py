#!/usr/bin/env python3
"""Generate W2 Agentic Architecture Patterns lecture slides (~65 slides)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import copy
import os

# ─── Load reference and strip slides ───
ref_path = "/Users/sungju/Downloads/Agentic AI - Intro.pptx"
prs = Presentation(ref_path)

# Remove all existing slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# ─── Style constants (from reference analysis) ───
FONT_TITLE = "Arial"
FONT_BODY = "Inter"
FONT_BODY_LIGHT = "Inter Light"
FONT_BODY_BOLD = "Inter"
COLOR_BODY = RGBColor(0x43, 0x43, 0x43)
COLOR_REF = RGBColor(0x66, 0x66, 0x66)
COLOR_ACCENT = RGBColor(0x42, 0x85, 0xF4)  # Google blue
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BLACK = RGBColor(0x00, 0x00, 0x00)
COLOR_DK2 = RGBColor(0x59, 0x59, 0x59)

FIG_DIR = "figures/cropped"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# ─── Helper functions ───

def add_title_slide(title_text, subtitle_text=""):
    """Add a title slide (layout 0)."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    # Title placeholder idx=0
    title_ph = slide.placeholders[0]
    title_ph.text = title_text
    for para in title_ph.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = COLOR_BLACK
    # Subtitle placeholder idx=1
    if subtitle_text and 1 in slide.placeholders:
        sub_ph = slide.placeholders[1]
        sub_ph.text = subtitle_text
        for para in sub_ph.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = COLOR_DK2
    return slide

def add_section_slide(title_text):
    """Add a section header slide (layout 1)."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_ph = slide.placeholders[0]
    title_ph.text = title_text
    for para in title_ph.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = COLOR_BLACK
    return slide

def add_content_slide(title_text, body_items, reference="", figure_path=None,
                       figure_pos=None, body_top=None, body_width=None,
                       bullet_size=Pt(11)):
    """Add a TITLE_ONLY slide with custom text boxes and optional figure.
    
    body_items: list of (text, is_bold, indent_level) tuples
    """
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    
    # Set title
    title_ph = slide.placeholders[0]
    title_ph.text = title_text
    for para in title_ph.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = COLOR_BLACK
    
    # Determine layout based on figure presence
    if figure_path and os.path.exists(figure_path):
        # Figure on right, text on left
        txt_left = Inches(0.4)
        txt_top = body_top or Inches(1.2)
        txt_width = body_width or Inches(4.8)
        txt_height = Inches(3.6)
        
        # Add figure
        if figure_pos:
            fig_left, fig_top, fig_width, fig_height = figure_pos
        else:
            fig_left = Inches(5.4)
            fig_top = Inches(1.2)
            fig_width = Inches(4.2)
            fig_height = Inches(3.2)
        
        try:
            slide.shapes.add_picture(figure_path, fig_left, fig_top, fig_width, fig_height)
        except:
            pass
    else:
        txt_left = Inches(0.4)
        txt_top = body_top or Inches(1.2)
        txt_width = body_width or Inches(9.0)
        txt_height = Inches(3.8)
    
    # Add body text box
    txBox = slide.shapes.add_textbox(txt_left, txt_top, txt_width, txt_height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(body_items):
        if isinstance(item, tuple):
            text, is_bold, indent = item
        else:
            text, is_bold, indent = item, False, 0
        
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        
        run = para.add_run()
        run.text = text
        run.font.size = bullet_size
        run.font.name = FONT_BODY if is_bold else FONT_BODY_LIGHT
        run.font.bold = is_bold
        run.font.color.rgb = COLOR_BODY
        
        para.level = indent
        para.space_after = Pt(4)
        
        if indent == 0 and not is_bold:
            # Regular bullet
            pass
    
    # Add reference at bottom
    if reference:
        ref_box = slide.shapes.add_textbox(Inches(0.4), Inches(5.05), Inches(9.0), Inches(0.4))
        ref_tf = ref_box.text_frame
        ref_tf.word_wrap = True
        ref_para = ref_tf.paragraphs[0]
        ref_run = ref_para.add_run()
        ref_run.text = reference
        ref_run.font.size = Pt(7)
        ref_run.font.name = FONT_BODY_LIGHT
        ref_run.font.color.rgb = COLOR_REF
    
    return slide

def add_figure_slide(title_text, figure_path, reference="", caption="",
                      fig_top=None, fig_height=None):
    """Add a slide with a large centered figure."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    
    # Set title
    title_ph = slide.placeholders[0]
    title_ph.text = title_text
    for para in title_ph.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = COLOR_BLACK
    
    if figure_path and os.path.exists(figure_path):
        from PIL import Image
        img = Image.open(figure_path)
        img_w, img_h = img.size
        aspect = img_w / img_h
        
        max_w = Inches(8.5)
        max_h = fig_height or Inches(3.5)
        
        # Calculate dimensions preserving aspect ratio
        if aspect > (max_w / max_h):
            w = max_w
            h = int(w / aspect)
        else:
            h = max_h
            w = int(h * aspect)
        
        left = int((SLIDE_W - w) / 2)
        top = fig_top or Inches(1.3)
        
        try:
            slide.shapes.add_picture(figure_path, left, top, w, h)
        except:
            pass
    
    # Add caption if provided
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.7), Inches(9.0), Inches(0.3))
        cap_tf = cap_box.text_frame
        cap_para = cap_tf.paragraphs[0]
        cap_run = cap_para.add_run()
        cap_run.text = caption
        cap_run.font.size = Pt(9)
        cap_run.font.name = FONT_BODY_LIGHT
        cap_run.font.color.rgb = COLOR_DK2
        cap_para.alignment = PP_ALIGN.CENTER
    
    # Add reference
    if reference:
        ref_box = slide.shapes.add_textbox(Inches(0.4), Inches(5.05), Inches(9.0), Inches(0.4))
        ref_tf = ref_box.text_frame
        ref_para = ref_tf.paragraphs[0]
        ref_run = ref_para.add_run()
        ref_run.text = reference
        ref_run.font.size = Pt(7)
        ref_run.font.name = FONT_BODY_LIGHT
        ref_run.font.color.rgb = COLOR_REF
    
    return slide

def add_text_slide(title_text, body_items, reference="", bullet_size=Pt(11)):
    """Add a full-width text slide using TITLE_AND_BODY layout (layout 2)."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    
    # Set title
    title_ph = slide.placeholders[0]
    title_ph.text = title_text
    for para in title_ph.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = COLOR_BLACK
    
    # Add body text box
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(1.2), Inches(9.0), Inches(3.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(body_items):
        if isinstance(item, tuple):
            if len(item) == 3:
                text, is_bold, indent = item
            else:
                text, is_bold = item
                indent = 0
        else:
            text, is_bold, indent = item, False, 0
        
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        
        run = para.add_run()
        run.text = text
        run.font.size = bullet_size
        run.font.name = FONT_BODY if is_bold else FONT_BODY_LIGHT
        run.font.bold = is_bold
        run.font.color.rgb = COLOR_BODY
        para.level = indent
        para.space_after = Pt(4)
    
    if reference:
        ref_box = slide.shapes.add_textbox(Inches(0.4), Inches(5.05), Inches(9.0), Inches(0.4))
        ref_tf = ref_box.text_frame
        ref_para = ref_tf.paragraphs[0]
        ref_run = ref_para.add_run()
        ref_run.text = reference
        ref_run.font.size = Pt(7)
        ref_run.font.name = FONT_BODY_LIGHT
        ref_run.font.color.rgb = COLOR_REF
    
    return slide

# ════════════════════════════════════════════════════════════════
# SLIDE GENERATION
# ════════════════════════════════════════════════════════════════

print("Creating slides...")

# ─── SLIDE 1: Title ───
add_title_slide(
    "Agentic Architecture Patterns",
    "W2 \u2014 Agentic AI (Spring 2026)"
)

# ─── SLIDE 2: Agenda ───
add_text_slide("Lecture Outline", [
    ("Part 1: Chain-of-Thought & Reasoning", True, 0),
    ("CoT Prompting, Self-Consistency, Least-to-Most, RAP, Test-Time Compute Scaling", False, 1),
    ("", False, 0),
    ("Part 2: Search & Planning", True, 0),
    ("Tree of Thoughts, Plan-and-Solve, LATS, Graph of Thoughts, Planning Survey, Planning Abilities", False, 1),
    ("", False, 0),
    ("Part 3: Grounded Acting & ReAct", True, 0),
    ("ReAct, Reflexion, FireAct, Scaling Agent Systems", False, 1),
])

# ════════════════════════════════════════════════════════════════
# SECTION 1: Chain-of-Thought & Reasoning
# ════════════════════════════════════════════════════════════════
add_section_slide("Part 1: Chain-of-Thought & Reasoning")

# ─── PAPER 1: Chain-of-Thought (4 slides) ───
ref_cot = "[Wei et al. 22] Chain-of-Thought Prompting Elicits Reasoning in Large Language Models, NeurIPS 2022"

# CoT Slide 1: Overview + Figure
add_content_slide(
    "Chain-of-Thought Prompting",
    [
        ("Core Idea: Augment few-shot exemplars with intermediate reasoning steps to enable multi-step reasoning", True, 0),
        ("", False, 0),
        ("Standard prompting provides (input, output) pairs \u2014 CoT provides (input, chain-of-thought, output) triples", False, 0),
        ("", False, 0),
        ("No finetuning, no architectural changes \u2014 just add reasoning traces to prompt exemplars", False, 0),
        ("", False, 0),
        ("Emergent ability: only works with models \u2265100B parameters; smaller models produce fluent but illogical chains", False, 0),
    ],
    reference=ref_cot,
    figure_path=f"{FIG_DIR}/cot_fig1.png",
)

# CoT Slide 2: Method
add_text_slide(
    "CoT: How It Works",
    [
        ("Method: Chain-of-Thought Prompting", True, 0),
        ("Manually compose 8 exemplars with natural language reasoning steps", False, 1),
        ("At test time, model generates its own chain of thought mimicking demonstrated style", False, 1),
        ("Same 8 exemplars used across ALL math benchmarks (robust)", False, 1),
        ("", False, 0),
        ("Four Key Properties", True, 0),
        ("Decomposition: allocates more computation to harder problems (more steps)", False, 1),
        ("Interpretability: chain of thought provides window into model's reasoning", False, 1),
        ("Generality: applicable to math, commonsense, symbolic \u2014 any task humans solve via language", False, 1),
        ("Simplicity: just add examples of reasoning to the prompt", False, 1),
        ("", False, 0),
        ("Ablation: equation-only, variable-compute-only, and reasoning-after-answer all perform worse", True, 0),
    ],
    reference=ref_cot,
)

# CoT Slide 3: Results
add_content_slide(
    "CoT: Key Results",
    [
        ("GSM8K Math Reasoning", True, 0),
        ("PaLM 540B + CoT: 56.9% vs 17.9% standard prompting (3\u00d7 improvement)", False, 1),
        ("Surpasses finetuned GPT-3 175B + verifier (55%) \u2014 new SOTA", False, 1),
        ("", False, 0),
        ("Emergent Ability", True, 0),
        ("CoT does NOT help models <100B parameters", False, 1),
        ("Performance suddenly emerges at ~100B scale", False, 1),
        ("", False, 0),
        ("OOD Generalization", True, 0),
        ("Last Letter Concat: CoT enables generalization from 2-word to 3-4 word cases", False, 1),
        ("Standard prompting completely fails out-of-distribution", False, 1),
    ],
    reference=ref_cot,
    figure_path=f"{FIG_DIR}/cot_fig4.png",
)

# CoT Slide 4: Impact
add_text_slide(
    "CoT: Impact & Connections",
    [
        ("Foundation for ALL Modern Agent Reasoning", True, 0),
        ("Directly led to: Self-Consistency, Least-to-Most, Tree-of-Thought, Zero-shot CoT", False, 1),
        ("Key insight for agentic AI: models can plan and decompose if prompted correctly", False, 1),
        ("Inspired the 'System 2 thinking' paradigm in AI agents", False, 1),
        ("", False, 0),
        ("Limitations", True, 0),
        ("No guarantee of correct reasoning \u2014 can produce plausible but wrong reasoning paths", False, 1),
        ("Only works with very large models (~100B+) \u2014 costly for deployment", False, 1),
        ("Whether neural networks truly 'reason' vs pattern-match remains open", False, 1),
        ("", False, 0),
        ("Progression in This Lecture", True, 0),
        ("CoT (base) \u2192 Self-Consistency (better decoding) \u2192 Least-to-Most (harder problems)", False, 1),
    ],
    reference=ref_cot,
)

# ─── PAPER 2: Self-Consistency (3 slides) ───
ref_sc = "[Wang et al. 23] Self-Consistency Improves Chain of Thought Reasoning in Language Models, ICLR 2023"

add_content_slide(
    "Self-Consistency: Sample-and-Vote Decoding",
    [
        ("Core Idea: Sample multiple diverse reasoning paths, then majority vote on final answers", True, 0),
        ("", False, 0),
        ("Correct reasoning paths tend to converge on the same answer; incorrect paths scatter", False, 0),
        ("", False, 0),
        ("Three-Step Process", True, 0),
        ("1. Prompt with CoT (same exemplars as Wei et al.)", False, 1),
        ("2. Sample ~40 completions via temperature sampling", False, 1),
        ("3. Majority vote across all paths for final answer", False, 1),
        ("", False, 0),
        ("A 'self-ensemble' over a single model \u2014 no additional training or auxiliary models", False, 0),
    ],
    reference=ref_sc,
    figure_path=f"{FIG_DIR}/sc_fig1.png",
)

add_text_slide(
    "Self-Consistency: Results",
    [
        ("Arithmetic Reasoning (PaLM-540B)", True, 0),
        ("GSM8K: 56.5% \u2192 74.4% (+17.9% with 40 paths)", False, 1),
        ("MultiArith: 94.7% \u2192 99.3%", False, 1),
        ("SVAMP: 79.0% \u2192 86.6%", False, 1),
        ("AQuA: 35.8% \u2192 48.3%", False, 1),
        ("", False, 0),
        ("Key Findings", True, 0),
        ("Majority vote (unweighted) outperforms probability-weighted aggregation", False, 1),
        ("Performance improves consistently from 1 to 40 paths; saturates ~40", False, 1),
        ("Works even with zero-shot CoT: PaLM-540B GSM8K 43.0% \u2192 69.2% (+26.2%)", False, 1),
        ("Beam search yields LESS diversity \u2192 worse performance than sampling", False, 1),
        ("", False, 0),
        ("Consistency correlates with accuracy \u2192 serves as uncertainty estimate", True, 0),
    ],
    reference=ref_sc,
)

add_text_slide(
    "Self-Consistency: Impact",
    [
        ("Why It Matters", True, 0),
        ("Introduces the 'sample and aggregate' pattern for agent reliability", False, 1),
        ("Shows diversity in reasoning is a FEATURE, not a bug", False, 1),
        ("Foundation for majority voting / debate mechanisms in multi-agent systems", False, 1),
        ("", False, 0),
        ("Limitations", True, 0),
        ("Requires 10-40\u00d7 more computation (though parallelizable)", False, 1),
        ("Only works for fixed-answer problems with discrete answer spaces", False, 1),
        ("", False, 0),
        ("When CoT Hurts, Self-Consistency Recovers", True, 0),
        ("ANLI R1: Standard 69.1% \u2192 CoT 68.8% (hurts!) \u2192 Self-Consistency 78.5%", False, 1),
        ("Outperforms both standard and CoT prompting even when CoT alone degrades", False, 1),
    ],
    reference=ref_sc,
)

# ─── PAPER 3: Least-to-Most (4 slides) ───
ref_ltm = "[Zhou et al. 23] Least-to-Most Prompting Enables Complex Reasoning in LLMs, ICLR 2023"

add_content_slide(
    "Least-to-Most: Explicit Task Decomposition",
    [
        ("Core Idea: Decompose complex problems into simpler subproblems, solve sequentially from easiest to hardest", True, 0),
        ("", False, 0),
        ("CoT fails on problems harder than exemplars \u2014 Least-to-Most enables easy-to-hard generalization", False, 0),
        ("", False, 0),
        ("Two-Stage Approach", True, 0),
        ("Stage 1 (Decompose): LLM breaks question into ordered subquestions", False, 1),
        ("Stage 2 (Solve): Solve subproblems sequentially; each builds on previous answers", False, 1),
        ("", False, 0),
        ("Inspired by educational psychology: scaffold from simple to complex", False, 0),
    ],
    reference=ref_ltm,
    figure_path=f"{FIG_DIR}/ltm_fig1.png",
)

add_text_slide(
    "Least-to-Most: Dramatic Results on Compositional Generalization",
    [
        ("SCAN Length Split (the hardest split)", True, 0),
        ("Standard prompting: 16.7%", False, 1),
        ("Chain-of-Thought: 16.2% (CoT completely fails!)", False, 1),
        ("Least-to-Most: 99.7% with just 14 exemplars", False, 1),
        ("Previous neural-symbolic models needed 15,000+ training examples", False, 1),
        ("", False, 0),
        ("Last Letter Concatenation", True, 0),
        ("CoT degrades rapidly with length: L=4 84.2% \u2192 L=12 31.8%", False, 1),
        ("Least-to-Most maintains performance: L=4 94.0% \u2192 L=12 74.0%", False, 1),
        ("", False, 0),
        ("GSM8K (Math)", True, 0),
        ("Overall modest gain: 60.87% \u2192 62.39%", False, 1),
        ("But on hardest problems (5+ steps): 39.07% \u2192 45.23% (+6.16%)", False, 1),
        ("Shines on harder problems, slight underperformance on simple 2-step problems", False, 1),
    ],
    reference=ref_ltm,
)

add_text_slide(
    "Least-to-Most: Significance for Agentic AI",
    [
        ("Establishes Explicit Task Decomposition as a First-Class Strategy", True, 0),
        ("Core pattern adopted by: AutoGPT, BabyAGI, LangChain agents", False, 1),
        ("Directly inspired: Decomposed Prompting, Plan-and-Solve, Successive Prompting", False, 1),
        ("", False, 0),
        ("The Reasoning Prompting Stack", True, 0),
        ("1. CoT: 'Think step by step' (single path, single pass)", False, 1),
        ("2. Self-Consistency: 'Think multiple ways, vote' (multiple paths, single pass)", False, 1),
        ("3. Least-to-Most: 'Break it down, solve piece by piece' (single path, multiple passes)", False, 1),
        ("Combined: Decompose \u2192 solve each with CoT \u2192 self-consistency vote", False, 1),
        ("", False, 0),
        ("Limitation: Requires multiple sequential LLM calls (increased latency)", True, 0),
        ("Not all problems benefit from explicit decomposition", False, 1),
    ],
    reference=ref_ltm,
)

# ─── PAPER 4: RAP (4 slides) ───
ref_rap = "[Hao et al. 23] Reasoning with Language Model is Planning with World Model, EMNLP 2023"

add_content_slide(
    "RAP: LLM as World Model + Agent",
    [
        ("Core Idea: Reframe LLM reasoning as a planning problem using Monte Carlo Tree Search (MCTS)", True, 0),
        ("", False, 0),
        ("The LLM serves dual roles:", False, 0),
        ("World Model: simulates states and predicts outcomes of actions", False, 1),
        ("Reasoning Agent: generates candidate reasoning steps", False, 1),
        ("", False, 0),
        ("MCTS strategically explores reasoning space, balancing exploration vs exploitation", False, 0),
        ("", False, 0),
        ("Key finding: LLaMA-33B + RAP surpasses GPT-4 + CoT by 33% on Blocksworld", True, 0),
    ],
    reference=ref_rap,
    figure_path=f"{FIG_DIR}/rap_fig1.png",
)

add_figure_slide(
    "RAP: MCTS Planning Algorithm",
    f"{FIG_DIR}/rap_fig3.png",
    reference=ref_rap,
    caption="Four MCTS phases: Selection (UCT), Expansion (sample actions), Simulation (rollout), Backpropagation (update Q-values)",
)

add_text_slide(
    "RAP: Experimental Results",
    [
        ("Blocksworld (Plan Generation)", True, 0),
        ("LLaMA-33B: CoT 1% vs RAP 64% average success", False, 1),
        ("RAP(20): 2-step 100%, 4-step 88%, 6-step 42%", False, 1),
        ("LLaMA-33B + RAP surpasses GPT-4 + CoT", False, 1),
        ("", False, 0),
        ("GSM8K (Math Reasoning)", True, 0),
        ("CoT: 29.4% \u2192 CoT+SC(10): 46.8% \u2192 RAP+agg(10): 51.6%", False, 1),
        ("RAP with aggregation outperforms all baselines", False, 1),
        ("", False, 0),
        ("ProntoQA (Logical Reasoning)", True, 0),
        ("CoT proof accuracy: 64.8% \u2192 RAP: 78.8% (+14pp)", False, 1),
        ("", False, 0),
        ("Key Insight: Smaller models + better search can outperform larger models", True, 0),
        ("Directly foreshadows the test-time compute scaling thesis", False, 1),
    ],
    reference=ref_rap,
)

add_text_slide(
    "RAP: Reward Design & Impact",
    [
        ("Three Reward Types", True, 0),
        ("Action Likelihood: log probability reflects LLM's intuitive preference", False, 1),
        ("State Confidence: sample multiple answers, use proportion of most frequent", False, 1),
        ("Self-Evaluation: prompt LLM with 'Is this correct?' \u2192 P('Yes') as reward", False, 1),
        ("", False, 0),
        ("Connection to Other Methods", True, 0),
        ("RAP uses MCTS; ToT uses BFS/DFS \u2014 both explore tree-structured reasoning", False, 1),
        ("RAP's world model is more principled than ToT's heuristic evaluation", False, 1),
        ("Demonstrates search/planning at inference time dramatically improves reasoning", False, 1),
        ("", False, 0),
        ("Limitations", True, 0),
        ("High computational cost: many LLM calls per MCTS iteration", False, 1),
        ("Task-specific reward and state/action design still requires manual engineering", False, 1),
    ],
    reference=ref_rap,
)

# ─── PAPER 5: Test-Time Compute (4 slides) ───
ref_ttc = "[Snell et al. 25] Scaling LLM Test-Time Compute Optimally, ICML 2025"

add_content_slide(
    "Scaling Test-Time Compute",
    [
        ("Core Idea: How to optimally allocate additional computation at inference time to improve LLM performance", True, 0),
        ("", False, 0),
        ("Key Finding: Compute-optimal scaling improves efficiency by >4\u00d7 over naive best-of-N", False, 0),
        ("", False, 0),
        ("A smaller model with optimal test-time compute can outperform a 14\u00d7 larger model", True, 0),
        ("", False, 0),
        ("Effectiveness depends critically on problem difficulty \u2192 need adaptive strategy selection", False, 0),
    ],
    reference=ref_ttc,
    figure_path=f"{FIG_DIR}/ttc_fig1.png",
)

add_figure_slide(
    "Test-Time Compute: Two Strategies",
    f"{FIG_DIR}/ttc_fig2.png",
    reference=ref_ttc,
    caption="Best-of-N (sample N solutions), Beam Search (prune at each step), Lookahead Search (roll out k steps ahead)",
)

add_text_slide(
    "Test-Time Compute: Key Results",
    [
        ("Two Knobs for Test-Time Compute", True, 0),
        ("1. Verifier-guided search: Best-of-N, Beam Search, Lookahead Search", False, 1),
        ("2. Iterative revision: model conditions on previous (incorrect) attempts", False, 1),
        ("", False, 0),
        ("Compute-Optimal Strategy", True, 0),
        ("Estimate problem difficulty using base model's pass@1 rate", False, 1),
        ("Easy questions: sequential revisions work best", False, 1),
        ("Hard questions: balanced mix of sequential + parallel sampling", False, 1),
        ("Hardest questions: no method makes meaningful progress \u2014 need better pretraining", False, 1),
        ("", False, 0),
        ("Test-Time vs Pretraining Compute", True, 0),
        ("NOT 1-to-1 exchangeable \u2014 depends on difficulty and inference load", False, 1),
        ("Easy/medium questions: test-time compute is preferable to scaling parameters", False, 1),
        ("Hardest questions: pretraining is almost always more effective", False, 1),
    ],
    reference=ref_ttc,
)

add_text_slide(
    "Test-Time Compute: Implications",
    [
        ("Three Key Takeaways", True, 0),
        ("Verifiers: Beam search better for harder questions at lower budgets; Best-of-N better for easy questions", False, 1),
        ("Revisions: Ideal sequential-to-parallel ratio depends on difficulty", False, 1),
        ("Scaling: test-time compute + pretraining are complementary, not interchangeable", False, 1),
        ("", False, 0),
        ("Impact", True, 0),
        ("Directly motivated OpenAI o1, DeepSeek-R1 \u2014 heavy test-time computation", False, 1),
        ("New scaling law dimension: alongside data and parameters, now inference compute", False, 1),
        ("Future: smaller models + more test-time compute replacing large models", False, 1),
        ("", False, 0),
        ("Connection to This Lecture", True, 0),
        ("RAP/ToT are specific instantiations of test-time compute strategies", False, 1),
        ("This paper provides the meta-analysis of WHEN each strategy works best", False, 1),
        ("Validates: smaller model + search can close up to 14\u00d7 model size gap", False, 1),
    ],
    reference=ref_ttc,
)

# ════════════════════════════════════════════════════════════════
# SECTION 2: Search & Planning
# ════════════════════════════════════════════════════════════════
add_section_slide("Part 2: Search & Planning")

# ─── PAPER 6: Tree of Thoughts (4 slides) ───
ref_tot = "[Yao et al. 23] Tree of Thoughts: Deliberate Problem Solving with LLMs, NeurIPS 2023"

add_content_slide(
    "Tree of Thoughts (ToT)",
    [
        ("Core Idea: Generalize CoT from linear chain to tree-structured search with BFS/DFS", True, 0),
        ("", False, 0),
        ("LLM generates multiple reasoning paths ('thoughts'), self-evaluates intermediate states, and uses search algorithms with lookahead and backtracking", False, 0),
        ("", False, 0),
        ("Inspired by dual-process theory: augment System 1 (fast, automatic) with System 2 (slow, deliberate)", False, 0),
        ("", False, 0),
        ("Game of 24: CoT achieves 4%, ToT achieves 74% \u2014 18.5\u00d7 improvement", True, 0),
    ],
    reference=ref_tot,
    figure_path=f"{FIG_DIR}/tot_fig1.png",
)

add_text_slide(
    "ToT: Four Design Decisions",
    [
        ("1. Thought Decomposition", True, 0),
        ("Task-dependent granularity: words (Game of 24), lines (Crosswords), paragraphs (Writing)", False, 1),
        ("Must be 'small enough' for diverse generation, 'big enough' to evaluate progress", False, 1),
        ("", False, 0),
        ("2. Thought Generator G(p\u03b8, s, k)", True, 0),
        ("(a) Sample i.i.d. from CoT prompt (rich thought spaces)", False, 1),
        ("(b) Propose sequentially in one pass (constrained thought spaces)", False, 1),
        ("", False, 0),
        ("3. State Evaluator V(p\u03b8, S)", True, 0),
        ("(a) Value each state independently: 'sure/maybe/impossible'", False, 1),
        ("(b) Vote across states: LM deliberates and picks most promising", False, 1),
        ("", False, 0),
        ("4. Search Algorithm", True, 0),
        ("BFS for shallow trees (depth 2-3): Game of 24, Creative Writing", False, 1),
        ("DFS with pruning for deeper search (10+ steps): Crosswords", False, 1),
    ],
    reference=ref_tot,
)

add_text_slide(
    "ToT: Results Across Three Tasks",
    [
        ("Game of 24 (arithmetic puzzle)", True, 0),
        ("IO: 7.3%, CoT: 4.0%, CoT-SC(k=100): 9.0%, ToT(b=5): 74%", False, 1),
        ("~60% of CoT failures happen at the very first step", False, 1),
        ("", False, 0),
        ("Creative Writing (coherent 4-paragraph essay)", True, 0),
        ("GPT-4 coherency: IO 6.19, CoT 6.93, ToT 7.56", False, 1),
        ("Humans prefer ToT over CoT in 41/100 cases (vs 21 for CoT)", False, 1),
        ("", False, 0),
        ("Mini Crosswords (5\u00d75)", True, 0),
        ("Word success: IO 14%, CoT 15.6%, ToT 60% (4\u00d7 improvement)", False, 1),
        ("Removing backtracking drops success from 60% to 20%", False, 1),
        ("", False, 0),
        ("Key: IO, CoT, CoT-SC, self-refinement are all special cases of ToT (limited depth/breadth)", True, 0),
    ],
    reference=ref_tot,
)

add_text_slide(
    "ToT: Impact & Connections",
    [
        ("One of the Most Influential Prompting Papers", True, 0),
        ("Established 'tree search over thoughts' as standard paradigm", False, 1),
        ("Directly inspired: RAP (MCTS), Graph of Thoughts, Algorithm of Thoughts", False, 1),
        ("Influenced reasoning agents and test-time compute scaling literature", False, 1),
        ("", False, 0),
        ("Limitations", True, 0),
        ("Requires task-specific prompt engineering for decomposition & evaluation", False, 1),
        ("ToT with b=5 requires ~15\u00d7 more LLM calls than CoT", False, 1),
        ("Only tested on GPT-4; only uses BFS/DFS (not A* or MCTS)", False, 1),
        ("", False, 0),
        ("Connection to Other Papers", True, 0),
        ("ToT (paradigm) \u2192 RAP (principled MCTS) \u2192 Test-Time Compute (optimal allocation)", False, 1),
        ("ToT uses LM as evaluator via prompting; RAP formalizes as reward function", False, 1),
    ],
    reference=ref_tot,
)

# ─── PAPER 7: Plan-and-Solve (3 slides) ───
ref_ps = "[Wang et al. 23] Plan-and-Solve Prompting, ACL 2023"

add_content_slide(
    "Plan-and-Solve: Improving Zero-Shot CoT",
    [
        ("Problem: Zero-shot CoT ('Let's think step by step') has three systematic errors", True, 0),
        ("Calculation errors (7%), missing steps (12%), semantic misunderstanding (27%)", False, 1),
        ("", False, 0),
        ("Solution: Replace vague trigger with explicit planning instruction", True, 0),
        ("'Let's first understand the problem and devise a plan to solve it. Then, let's carry out the plan step by step.'", False, 1),
        ("", False, 0),
        ("PS+ adds: 'extract relevant variables', 'calculate intermediate results', 'pay attention to calculation'", False, 0),
    ],
    reference=ref_ps,
    figure_path=f"{FIG_DIR}/ps_fig2.png",
)

add_text_slide(
    "Plan-and-Solve: Results",
    [
        ("Zero-Shot PS+ vs Zero-Shot CoT (GPT-3, text-davinci-003)", True, 0),
        ("MultiArith: 83.8% \u2192 91.8% (+8.0%)", False, 1),
        ("AddSub: 85.3% \u2192 92.2% (+6.9%)", False, 1),
        ("GSM8K: 56.4% \u2192 59.3% (+2.9%)", False, 1),
        ("AQuA: 38.9% \u2192 46.0% (+7.1%)", False, 1),
        ("Average: 70.4% \u2192 76.7% (+6.3%)", False, 1),
        ("", False, 0),
        ("Remarkable: Zero-shot PS+ (76.7%) nearly matches 8-shot Manual-CoT (77.6%)", True, 0),
        ("", False, 0),
        ("With Self-Consistency: GSM8K 73.7%, SVAMP 84.4%", False, 0),
        ("PS+ reduces calculation errors from 7% to 5%, missing-step errors from 12% to 7%", False, 0),
    ],
    reference=ref_ps,
)

add_text_slide(
    "Plan-and-Solve: Significance",
    [
        ("Establishes 'Plan Then Execute' Paradigm", True, 0),
        ("Central to agentic AI: separate planning and action phases", False, 1),
        ("Shows prompt engineering can achieve near few-shot performance at zero cost", False, 1),
        ("", False, 0),
        ("90 out of 100 sampled predictions indeed incorporated an explicit plan", False, 0),
        ("", False, 0),
        ("Limitation: Semantic misunderstanding errors (27%) remain \u2014 these stem from model capability, not prompting", True, 0),
        ("", False, 0),
        ("Connection: PS is a single-chain approach (linear plan then linear execution)", False, 0),
        ("GoT generalizes this to arbitrary graph structures", False, 1),
        ("LATS adds tree search with backtracking and external feedback", False, 1),
    ],
    reference=ref_ps,
)

# ─── PAPER 8: LATS (4 slides) ───
ref_lats = "[Zhou et al. 24] LATS: Language Agent Tree Search Unifies Reasoning, Acting, and Planning, ICML 2024"

add_content_slide(
    "LATS: Unifying Reasoning, Acting, and Planning",
    [
        ("Core Idea: Adapt Monte Carlo Tree Search to language agents that interact with environments", True, 0),
        ("", False, 0),
        ("LLM serves triple duty: agent (actions), value function (evaluating states), and reflector (self-critique)", False, 0),
        ("", False, 0),
        ("Unlike prior methods: LATS can backtrack, explore alternatives, and learn from failures through reflection", False, 0),
        ("", False, 0),
        ("The ONLY method with all five: Reasoning, Acting, Planning, Self-Reflection, External Memory", True, 0),
    ],
    reference=ref_lats,
    figure_path=f"{FIG_DIR}/lats_fig1.png",
)

add_figure_slide(
    "LATS: Six Operations in the Search Loop",
    f"{FIG_DIR}/lats_fig2.png",
    reference=ref_lats,
    caption="Selection (UCT) \u2192 Expansion (sample actions) \u2192 Evaluation (LM score + self-consistency) \u2192 Simulation \u2192 Backpropagation \u2192 Reflection",
)

add_text_slide(
    "LATS: State-of-the-Art Results",
    [
        ("HotPotQA (GPT-3.5)", True, 0),
        ("LATS: 0.71 EM vs Reflexion 0.51, ReAct 0.32 \u2014 doubles ReAct", False, 1),
        ("", False, 0),
        ("HumanEval (Code Generation)", True, 0),
        ("GPT-3.5: LATS 83.8 Pass@1 vs Reflexion 68.1, ReAct 56.9", False, 1),
        ("GPT-4: LATS 92.7 Pass@1 vs Reflexion 91.0 \u2014 new SOTA", False, 1),
        ("", False, 0),
        ("WebShop (Web Navigation)", True, 0),
        ("LATS Score 75.9, SR 38.0 vs Reflexion 64.2/35.0", False, 1),
        ("Surpasses RL-based training (IL+RL: 62.4/28.7)", False, 1),
        ("", False, 0),
        ("Ablation: Removing LM heuristic causes largest drop (0.63 \u2192 0.37)", True, 0),
        ("MCTS vs DFS: MCTS outperforms (0.63 vs 0.42)", False, 1),
        ("Requires 3.55\u00d7 fewer nodes than RAP, 12.12\u00d7 fewer than ToT", False, 1),
    ],
    reference=ref_lats,
)

add_text_slide(
    "LATS: Evaluation + Value Function Design",
    [
        ("Novel Dual-Component Value Function", True, 0),
        ("V(s) = \u03bb \u00b7 LM(s) + (1-\u03bb) \u00b7 SC(s)", False, 1),
        ("LM(s): self-generated LM score (reasoning about trajectory correctness)", False, 1),
        ("SC(s): self-consistency score (actions sampled multiple times = more accurate)", False, 1),
        ("Evaluated AFTER receiving environmental feedback (unlike ToT)", False, 1),
        ("", False, 0),
        ("Reflection Mechanism", True, 0),
        ("Upon failure: LLM generates verbal self-reflection summarizing errors", False, 1),
        ("Failed trajectories + reflections stored in memory for future iterations", False, 1),
        ("", False, 0),
        ("Connection", True, 0),
        ("LATS extends ReAct's action loop with MCTS-based planning", False, 1),
        ("Bridges ToT's pure reasoning search with real environment interaction", False, 1),
        ("Limitation: assumes ability to revert to earlier states (may not hold in real-world)", False, 1),
    ],
    reference=ref_lats,
)

# ─── PAPER 9: Graph of Thoughts (4 slides) ───
ref_got = "[Besta et al. 24] Graph of Thoughts: Solving Elaborate Problems with LLMs, AAAI 2024"

add_content_slide(
    "Graph of Thoughts (GoT)",
    [
        ("Core Idea: Model LLM reasoning as an arbitrary directed graph, not just chains or trees", True, 0),
        ("", False, 0),
        ("Subsumes all prior paradigms: CoT (chain), CoT-SC (parallel chains), ToT (tree)", False, 0),
        ("", False, 0),
        ("Novel graph transformations:", True, 0),
        ("Aggregation: merge multiple thoughts (e.g., combine sorted subarrays)", False, 1),
        ("Refinement: loop over a thought iteratively (self-improvement)", False, 1),
        ("Generation: branch from existing thought (like ToT expansion)", False, 1),
    ],
    reference=ref_got,
    figure_path=f"{FIG_DIR}/got_fig1.png",
)

add_text_slide(
    "GoT: System Architecture",
    [
        ("GoT = (G, T, E, R): Graph, Transformations, Evaluator, Ranking", True, 0),
        ("", False, 0),
        ("Four Modules", True, 0),
        ("Prompter: prepares messages encoding graph structure + instructions", False, 1),
        ("Parser: extracts info from LLM responses to update thought state", False, 1),
        ("Scoring/Validation: evaluates thought quality (LLM, humans, or simple functions)", False, 1),
        ("Controller: orchestrates reasoning with Graph of Operations (GoO) + Graph Reasoning State (GRS)", False, 1),
        ("", False, 0),
        ("Key Design: Graph of Operations (GoO) is STATIC, defined upfront", True, 0),
        ("Like a program specifying which transformations to apply in what order", False, 1),
        ("Provides predictability and cost control", False, 1),
        ("Requires developer to design decomposition strategy per task", False, 1),
    ],
    reference=ref_got,
)

add_text_slide(
    "GoT: Results",
    [
        ("Sorting (GPT-3.5, 100 samples)", True, 0),
        ("32 elements: GoT ~2 errors vs ToT ~6, CoT ~9", False, 1),
        ("64 elements: GoT ~4 vs ToT ~28, CoT ~44", False, 1),
        ("128 elements: GoT ~8 vs ToT ~56, CoT ~88", False, 1),
        ("GoT reduces error by 62% vs ToT AND reduces cost by 31%", False, 1),
        ("", False, 0),
        ("Set Intersection: GoT 43/100 correct (32 elem) vs ToT 29, CoT 6", False, 0),
        ("Keyword Counting: GoT ~2 median errors vs ToT ~8", False, 0),
        ("", False, 0),
        ("Latency-Volume Tradeoff (Theoretical)", True, 0),
        ("CoT: Latency O(N), Volume O(N)", False, 1),
        ("ToT: Latency O(log N), Volume O(log N)", False, 1),
        ("GoT: Latency O(log N), Volume O(N) \u2014 both low latency AND high volume!", False, 1),
        ("", False, 0),
        ("GoT advantages increase with problem complexity", True, 0),
    ],
    reference=ref_got,
)

add_text_slide(
    "GoT: Impact & Evolution of Reasoning Topology",
    [
        ("Evolution of LLM Reasoning Structures", True, 0),
        ("IO (single node) \u2192 CoT (chain) \u2192 CoT-SC (parallel chains) \u2192 ToT (tree) \u2192 GoT (arbitrary graph)", False, 1),
        ("", False, 0),
        ("GoT's Unique Contribution: Aggregation", True, 0),
        ("Neither chains nor trees can express merging results from different branches", False, 1),
        ("Enables divide-and-conquer: decompose, solve independently, merge solutions", False, 1),
        ("", False, 0),
        ("Comparison with LATS", True, 0),
        ("GoT: static reasoning graph, no environment interaction, cost-efficient", False, 1),
        ("LATS: dynamic tree search, environment feedback, higher cost", False, 1),
        ("GoT excels at decomposable reasoning; LATS excels at interactive decision-making", False, 1),
        ("", False, 0),
        ("Limitation: GoO is static \u2014 cannot adapt decomposition strategy mid-execution", True, 0),
    ],
    reference=ref_got,
)

# ─── PAPER 10: Planning Survey (3 slides) ───
ref_psur = "[Huang et al. 24] Understanding the Planning of LLM Agents: A Survey, arXiv 2024"

add_content_slide(
    "LLM Agent Planning: A Comprehensive Taxonomy",
    [
        ("First systematic taxonomy of LLM-based agent planning", True, 0),
        ("", False, 0),
        ("Five Directions for Planning:", True, 0),
        ("1. Task Decomposition (divide and conquer)", False, 1),
        ("2. Multi-Plan Selection (generate and select)", False, 1),
        ("3. External Planner-Aided Planning", False, 1),
        ("4. Reflection & Refinement", False, 1),
        ("5. Memory-Augmented Planning", False, 1),
        ("", False, 0),
        ("Key Insight: methods are interconnected, not mutually exclusive \u2014 real agents combine strategies", True, 0),
    ],
    reference=ref_psur,
    figure_path=f"{FIG_DIR}/psur_fig1.png",
)

add_text_slide(
    "Planning Taxonomy: Five Strategies in Detail",
    [
        ("1. Task Decomposition", True, 0),
        ("Decomposition-First (HuggingGPT, Plan-and-Solve) vs Interleaved (CoT, ReAct)", False, 1),
        ("", False, 0),
        ("2. Multi-Plan Selection", True, 0),
        ("Self-Consistency (majority vote), ToT (BFS/DFS), GoT, MCTS, A* search", False, 1),
        ("", False, 0),
        ("3. External Planner-Aided", True, 0),
        ("LLM formalizes task; symbolic planner generates correct plan (LLM+P, PDDL)", False, 1),
        ("", False, 0),
        ("4. Reflection & Refinement", True, 0),
        ("Generate plan \u2192 reflect on failures \u2192 refine iteratively (Reflexion, CRITIC)", False, 1),
        ("", False, 0),
        ("5. Memory-Augmented", True, 0),
        ("Store past experiences as text (Generative Agents) or embeddings (TiM)", False, 1),
    ],
    reference=ref_psur,
    bullet_size=Pt(10),
)

add_text_slide(
    "Planning Survey: Benchmark Evaluation",
    [
        ("Six Methods Evaluated on 4 Benchmarks (text-davinci-003)", True, 0),
        ("", False, 0),
        ("ALFWorld: Reflexion 71% > CoT-SC 57% > ReAct 57% > Few-CoT 43%", False, 0),
        ("HotPotQA: Reflexion 39% > ReAct 34% > CoT-SC 33%", False, 0),
        ("FEVER: Reflexion 68% > ReAct 63% > CoT-SC 62%", False, 0),
        ("", False, 0),
        ("Three Key Findings", True, 0),
        ("Performance increases with token expenditure \u2014 Reflexion best but most expensive", False, 1),
        ("Few-shot examples are necessary for complex tasks", False, 1),
        ("Reflection is crucial: Reflexion uses ~2\u00d7 tokens vs ReAct but shows clear improvement", False, 1),
        ("", False, 0),
        ("Open Challenges: hallucination, plan feasibility, efficiency, multi-modal feedback", True, 0),
    ],
    reference=ref_psur,
)

# ─── PAPER 11: Planning Abilities (3 slides) ───
ref_pa = "[Valmeekam et al. 23] On the Planning Abilities of LLMs \u2014 A Critical Investigation, NeurIPS 2023"

add_content_slide(
    "Can LLMs Really Plan? A Critical Investigation",
    [
        ("Core Finding: LLMs achieve only ~3% success rate on autonomous plan generation", True, 0),
        ("", False, 0),
        ("Even for simple Blocksworld tasks that humans solve at 78%", False, 0),
        ("", False, 0),
        ("Seven test cases: plan generation, optimal planning, replanning, generalization, robustness", False, 0),
        ("", False, 0),
        ("BUT: LLM plans can seed formal planners \u2014 ALL plans repaired within 7 seconds", True, 0),
    ],
    reference=ref_pa,
    figure_path=f"{FIG_DIR}/pa_fig1.png",
)

add_text_slide(
    "Planning Abilities: Sobering Results",
    [
        ("Autonomous Mode (Blocksworld, 600 instances)", True, 0),
        ("Plan Generation: GPT-3 1%, InstructGPT-3 6.8%, BLOOM 1.6%", False, 1),
        ("Optimal Planning: GPT-3 0.3%, InstructGPT-3 5.8%", False, 1),
        ("Neither model could solve ANY 5-block instance", False, 1),
        ("", False, 0),
        ("Goal Reformulation scores are misleading (77.8%)", True, 0),
        ("LLM just repeats the example plan \u2014 pattern matching, not reasoning", False, 1),
        ("When example plan is swapped, accuracy drops drastically", False, 1),
        ("", False, 0),
        ("Mystery Blocksworld (disguised domain names)", True, 0),
        ("InstructGPT-3 drops to 1.1% with deceptive names", False, 1),
        ("Proves LLMs rely on semantic meaning of actions, not abstract reasoning", False, 1),
        ("", False, 0),
        ("Fine-tuning (1000 instances): only improves to 16.4% \u2014 still very low", False, 0),
    ],
    reference=ref_pa,
)

add_text_slide(
    "Planning Abilities: Implications for Agent Architecture",
    [
        ("LLMs as Heuristic Guides, Not Autonomous Planners", True, 0),
        ("LLM plans are 'in the neighborhood' of correct (Levenshtein distance ~7 vs plan length ~11)", False, 1),
        ("LPG planner repairs ALL LLM plans to valid solutions within 7 seconds", False, 1),
        ("", False, 0),
        ("Architectural Implication", True, 0),
        ("LLM for reasoning/NL understanding + external tools/planners for correctness", False, 1),
        ("Motivates 'External Planner-Aided' category in the planning survey", False, 1),
        ("", False, 0),
        ("Key Tension in This Lecture", True, 0),
        ("ReAct (2023): 'LLMs plan effectively when grounded in environments'", False, 1),
        ("Valmeekam (2023): 'LLMs fundamentally cannot plan autonomously'", False, 1),
        ("Resolution: ReAct works because environment feedback corrects LLM errors \u2014 not because LLM truly plans", False, 1),
        ("", False, 0),
        ("Modern agents: hybrid architectures combining LLM reasoning with external verification", True, 0),
    ],
    reference=ref_pa,
)

# ════════════════════════════════════════════════════════════════
# SECTION 3: Grounded Acting & ReAct
# ════════════════════════════════════════════════════════════════
add_section_slide("Part 3: Grounded Acting & ReAct")

# ─── PAPER 12: ReAct (4 slides) ───
ref_react = "[Yao et al. 23] ReAct: Synergizing Reasoning and Acting in Language Models, ICLR 2023"

add_content_slide(
    "ReAct: Thought-Action-Observation Loop",
    [
        ("Core Idea: Interleave reasoning traces (Thoughts) and task-specific actions (Acts) in a loop", True, 0),
        ("", False, 0),
        ("CoT reasoning is ungrounded and hallucinates (56% of failures)", False, 0),
        ("Action-only agents fail to decompose goals or handle exceptions", False, 0),
        ("ReAct synergizes both: reasoning guides actions, observations inform reasoning", False, 0),
        ("", False, 0),
        ("Hallucination rate: CoT 56% vs ReAct 0%", True, 0),
    ],
    reference=ref_react,
    figure_path=f"{FIG_DIR}/react_fig1.png",
    figure_pos=(Inches(5.2), Inches(1.1), Inches(4.5), Inches(4.0)),
)

add_text_slide(
    "ReAct: Method",
    [
        ("Augmented Action Space: A\u0302 = A \u222a L", True, 0),
        ("Actions in A affect environment \u2192 produce Observations", False, 1),
        ("Actions in L are 'thoughts' (reasoning traces) \u2192 do NOT affect environment", False, 1),
        ("", False, 0),
        ("ReAct Loop", True, 0),
        ("1. Receive task/question and context", False, 1),
        ("2. Generate Thought (decompose, extract info, track progress)", False, 1),
        ("3. Generate Action (Search[entity], Lookup[string], Finish[answer])", False, 1),
        ("4. Receive Observation from environment", False, 1),
        ("5. Update context and repeat until completion", False, 1),
        ("", False, 0),
        ("Combining with Self-Consistency", True, 0),
        ("ReAct \u2192 CoT-SC: when ReAct fails, fall back to internal knowledge", False, 1),
        ("CoT-SC \u2192 ReAct: when CoT-SC has low confidence, fall back to search", False, 1),
        ("Best results: combining both modes", False, 1),
    ],
    reference=ref_react,
)

add_text_slide(
    "ReAct: Results Across Four Benchmarks",
    [
        ("Knowledge Tasks (PaLM-540B)", True, 0),
        ("HotPotQA: ReAct 27.4, CoT 29.4, Best combo (ReAct\u2192CoT-SC): 35.1 EM", False, 1),
        ("FEVER: ReAct 60.9 vs CoT 56.3 \u2014 ReAct outperforms", False, 1),
        ("Error: CoT 56% hallucination, ReAct 0% hallucination", False, 1),
        ("", False, 0),
        ("Decision Making", True, 0),
        ("ALFWorld: ReAct best trial 71% vs Act best 45% (+26% absolute)", False, 1),
        ("Even ReAct worst trial (48%) beats Act best (45%)", False, 1),
        ("WebShop: ReAct 40% SR vs Act 30.1% \u2014 outperforms IL+RL (28.7%)", False, 1),
        ("", False, 0),
        ("Fine-tuning Scaling (Figure 3)", True, 0),
        ("PaLM-8B finetuned ReAct outperforms ALL PaLM-62B prompting methods", False, 1),
        ("Fine-tuning teaches HOW to reason and act (generalizable skill)", False, 1),
    ],
    reference=ref_react,
)

add_text_slide(
    "ReAct: THE Foundational Paradigm for LLM Agents",
    [
        ("Impact: Directly Inspired Every Major Agent Framework", True, 0),
        ("LangChain agents, AutoGPT, BabyAGI all use Thought-Action loops", False, 1),
        ("Visual ChatGPT, HuggingGPT utilize ReAct's mechanism", False, 1),
        ("", False, 0),
        ("Unique Strengths", True, 0),
        ("Human-interpretable: reasoning traces can be inspected and verified", False, 1),
        ("Human-controllable: editing 1-2 thoughts changes agent behavior completely", False, 1),
        ("Zero hallucination in failure modes (vs 56% for CoT)", False, 1),
        ("General across: QA, fact verification, games, web navigation", False, 1),
        ("", False, 0),
        ("Limitations", True, 0),
        ("Interleaving constraint leads to more reasoning errors (47% vs CoT 16%)", False, 1),
        ("Non-informative search results derail reasoning (23% of errors)", False, 1),
        ("Prompting far from supervised SOTA (35.1 vs 67.5 HotPotQA)", False, 1),
    ],
    reference=ref_react,
)

# ─── PAPER 13: Reflexion (4 slides) ───
ref_refl = "[Shinn et al. 23] Reflexion: Language Agents with Verbal Reinforcement Learning, NeurIPS 2023"

add_content_slide(
    "Reflexion: Learning from Verbal Self-Critique",
    [
        ("Core Idea: Agents learn by generating natural language reflections on failures, stored in episodic memory", True, 0),
        ("", False, 0),
        ("Verbal reinforcement learning: converts sparse scalar/binary feedback into rich textual 'semantic gradients'", False, 0),
        ("", False, 0),
        ("No weight updates needed \u2014 the 'policy' is parameterized by {model, memory}", False, 0),
        ("", False, 0),
        ("AlfWorld: ReAct ~75% \u2192 ReAct + Reflexion 97% in 12 trials", True, 0),
    ],
    reference=ref_refl,
    figure_path=f"{FIG_DIR}/refl_fig2.png",
)

add_text_slide(
    "Reflexion: Three-Component Architecture",
    [
        ("1. Actor (M_a)", True, 0),
        ("LLM generating actions conditioned on state + memory", False, 1),
        ("Uses ReAct or CoT prompting as base agent framework", False, 1),
        ("", False, 0),
        ("2. Evaluator (M_e)", True, 0),
        ("Scores trajectory: exact-match (QA), heuristics (AlfWorld), self-generated unit tests (code)", False, 1),
        ("", False, 0),
        ("3. Self-Reflection Model (M_sr)", True, 0),
        ("Given sparse reward + trajectory + memory \u2192 generates nuanced verbal feedback", False, 1),
        ("Example: 'I incorrectly assumed both had the same profession... I should focus on accurately identifying...'", False, 1),
        ("", False, 0),
        ("Iterative Loop", True, 0),
        ("Trial 0: Act \u2192 Evaluate \u2192 Reflect \u2192 store reflection in memory", False, 1),
        ("Trial N: Act conditioned on memory \u2192 Evaluate \u2192 Reflect \u2192 append to memory", False, 1),
        ("Memory bounded to 1-3 reflections (context window limit)", False, 1),
    ],
    reference=ref_refl,
)

add_text_slide(
    "Reflexion: Results",
    [
        ("AlfWorld (134 tasks, sequential decision-making)", True, 0),
        ("ReAct only: 75%, plateaus with 22% hallucination", False, 1),
        ("ReAct + Reflexion: 97% (130/134 tasks) in 12 trials", False, 1),
        ("", False, 0),
        ("HotPotQA (reasoning)", True, 0),
        ("CoT + Reflexion: +14% accuracy improvement", False, 1),
        ("ReAct + Reflexion: +20% improvement", False, 1),
        ("Self-reflection adds +8% over episodic memory alone (ablation)", False, 1),
        ("", False, 0),
        ("Code Generation (HumanEval)", True, 0),
        ("GPT-4: 80.1% \u2192 91.0% (Reflexion) \u2014 new SOTA at time", False, 1),
        ("LeetcodeHard: 7.5% \u2192 15.0% (doubling performance)", False, 1),
        ("", False, 0),
        ("Both self-reflection AND test generation are necessary (ablation: neither alone improves)", True, 0),
    ],
    reference=ref_refl,
)

add_text_slide(
    "Reflexion: Impact on Agentic AI",
    [
        ("Foundational for Self-Improvement Paradigm", True, 0),
        ("Showed agents can learn from verbal experience \u2014 inspired memory-augmented agent designs", False, 1),
        ("Self-reflection more effective than blind retry", False, 1),
        ("Directly influenced FireAct (uses Reflexion trajectories as fine-tuning data)", False, 1),
        ("", False, 0),
        ("Limitations", True, 0),
        ("Natural language policy optimization may converge to local minima", False, 1),
        ("Long-term memory limited to sliding window; needs more advanced structures", False, 1),
        ("Self-generated unit tests can be flaky (16.3% false positive on MBPP)", False, 1),
        ("", False, 0),
        ("Connection to Scaling Agents paper", True, 0),
        ("Self-reflection is classified as single-agent mechanism (SAS)", False, 1),
        ("Stronger individual agents (from Reflexion) may reduce benefit of adding more agents", False, 1),
    ],
    reference=ref_refl,
)

# ─── PAPER 14: FireAct (3 slides) ───
ref_fa = "[Chen et al. 23] FireAct: Toward Language Agent Fine-tuning, arXiv 2023"

add_content_slide(
    "FireAct: Fine-tuning Language Agents",
    [
        ("Core Idea: Fine-tune small LMs on successful agent trajectories from strong teacher (GPT-4)", True, 0),
        ("", False, 0),
        ("Mix trajectories from multiple prompting methods (ReAct, CoT, Reflexion)", False, 0),
        ("", False, 0),
        ("Creates agents that implicitly select appropriate reasoning strategy at inference time", False, 0),
        ("", False, 0),
        ("Fine-tuned Llama-2-13B (34.4 EM) outperforms ALL GPT-3.5 prompting methods", True, 0),
    ],
    reference=ref_fa,
    figure_path=f"{FIG_DIR}/fa_fig2.png",
)

add_text_slide(
    "FireAct: Results & Efficiency",
    [
        ("Single-Task Fine-tuning (HotPotQA)", True, 0),
        ("Llama-2-7B: ReAct 14.8 \u2192 FireAct 26.2 EM (+77%)", False, 1),
        ("Llama-2-13B: 21.2 \u2192 34.4 EM (+62%)", False, 1),
        ("GPT-3.5: 31.4 \u2192 39.2 EM (+25%)", False, 1),
        ("", False, 0),
        ("Efficiency & Robustness", True, 0),
        ("Inference time: 9.0s \u2192 2.7s (70% reduction)", False, 1),
        ("Noisy tools: ReAct drops 33.8%, FireAct drops only 14.2%", False, 1),
        ("Generalization to unseen test set (Bamboogle): 44.0 vs ReAct 40.8", False, 1),
        ("", False, 0),
        ("Multi-Method Fine-tuning", True, 0),
        ("ReAct + CoT: 41.0 EM (best \u2014 fewer turns, faster)", False, 1),
        ("Oracle method selection: 52.0 EM \u2014 significant room for improvement", False, 1),
        ("Multi-task does not hurt source task; greatly improves target tasks", False, 1),
    ],
    reference=ref_fa,
)

add_text_slide(
    "FireAct: Bridging Prompting and Training",
    [
        ("Key Insight: Trajectory diversity > trajectory volume", True, 0),
        ("Mixing ReAct + CoT + Reflexion methods creates more capable agents", False, 1),
        ("500 trajectories sufficient for GPT-3.5; Llama models need 500-1000", False, 1),
        ("", False, 0),
        ("Distillation Paradigm", True, 0),
        ("Strong agent behavior from large models \u2192 distilled into smaller fine-tuned models", False, 1),
        ("Fine-tuning compiles Reflexion-style self-reflection into model weights", False, 1),
        ("", False, 0),
        ("Open Questions", True, 0),
        ("Multi-method agents don't yet reliably select best method (random 32.4 vs oracle 52.0)", False, 1),
        ("LoRA (26.2 EM) vs full fine-tuning (30.2 EM) \u2014 15% gap", False, 1),
        ("", False, 0),
        ("Connection: focuses on single-agent improvement through fine-tuning vs scaling agents' multi-agent approach", True, 0),
    ],
    reference=ref_fa,
)

# ─── PAPER 15: Scaling Agent Systems (4 slides) ───
ref_sa = "[Kim et al. 25] Towards a Science of Scaling Agent Systems, arXiv 2025 (Google Research)"

add_content_slide(
    "Scaling Agent Systems: More Agents \u2260 Better",
    [
        ("Core Finding: Multi-agent benefits are entirely task-contingent", True, 0),
        ("", False, 0),
        ("Finance Agent: +81% with centralized coordination", False, 0),
        ("PlanCraft: -70% degradation with independent agents", False, 0),
        ("", False, 0),
        ("Five Canonical Architectures Evaluated:", True, 0),
        ("SAS (Single-Agent), MAS-Independent, MAS-Decentralized, MAS-Centralized, MAS-Hybrid", False, 1),
        ("", False, 0),
        ("Predictive model: 87% accuracy in selecting optimal architecture", True, 0),
    ],
    reference=ref_sa,
    figure_path=f"{FIG_DIR}/sa_fig1.png",
)

add_text_slide(
    "Scaling Agents: Three Fundamental Principles",
    [
        ("1. Tool-Coordination Trade-off (\u03b2 = -0.267, p < 0.001)", True, 0),
        ("Tool-heavy tasks suffer from multi-agent overhead \u2014 MAS fragments per-agent token budget", False, 1),
        ("", False, 0),
        ("2. Capability Saturation (\u03b2 = -0.404, p < 0.001)", True, 0),
        ("Coordination yields diminishing/negative returns when SAS baselines exceed ~45%", False, 1),
        ("When single agent is already strong, adding agents hurts", False, 1),
        ("", False, 0),
        ("3. Topology-Dependent Error Amplification", True, 0),
        ("Independent agents amplify errors 17.2\u00d7 (unchecked propagation)", False, 1),
        ("Centralized coordination contains to 4.4\u00d7 (validation bottleneck intercepts errors)", False, 1),
        ("", False, 0),
        ("Overall MAS mean improvement: -3.5% (95% CI: [-18.6%, +25.7%])", True, 0),
        ("High variance (\u03c3 = 45.2%) confirms: task matters more than architecture", False, 0),
    ],
    reference=ref_sa,
)

add_text_slide(
    "Scaling Agents: When to Use Multi-Agent Systems",
    [
        ("Practical Architecture Selection Rules (87% accuracy)", True, 0),
        ("", False, 0),
        ("Parallelizable subtasks with independent information streams:", False, 0),
        ("\u2192 Centralized MAS (Finance Agent: +81%)", True, 1),
        ("", False, 0),
        ("High-entropy search spaces requiring exploration:", False, 0),
        ("\u2192 Decentralized MAS (BrowseComp: +9.2%)", True, 1),
        ("", False, 0),
        ("Sequential constraint-satisfaction tasks:", False, 0),
        ("\u2192 Single-Agent System (PlanCraft: MAS causes -70%)", True, 1),
        ("", False, 0),
        ("High baseline SAS accuracy (>45%):", False, 0),
        ("\u2192 Stay with SAS (capability saturation)", True, 1),
        ("", False, 0),
        ("Critical factor: task decomposability, NOT task complexity", True, 0),
    ],
    reference=ref_sa,
)

add_text_slide(
    "Scaling Agents: Validation & Connections",
    [
        ("Out-of-Sample Validation", True, 0),
        ("Tested on GPT-5.2 (released after study): MAE = 0.071", False, 1),
        ("4/5 scaling principles confirmed to generalize to unseen frontier models", False, 1),
        ("No single vendor achieves universal multi-agent dominance", False, 1),
        ("", False, 0),
        ("Connection to This Lecture", True, 0),
        ("Reflexion's self-reflection = single-agent mechanism (SAS)", False, 1),
        ("Stronger agents (from Reflexion/FireAct) may reduce benefit of adding more agents", False, 1),
        ("Validates: architectural choice matters as much as scaling model capability", False, 1),
        ("", False, 0),
        ("180 controlled configurations across 3 LLM families, 4 benchmarks", True, 0),
        ("First rigorous framework for agent system evaluation", False, 1),
    ],
    reference=ref_sa,
)

# ════════════════════════════════════════════════════════════════
# WRAP-UP SLIDES
# ════════════════════════════════════════════════════════════════

# Discussion Questions
add_text_slide(
    "Discussion Questions",
    [
        ("1. CoT, Self-Consistency, Least-to-Most form a 'reasoning prompting stack'. Can they be unified into a single framework? What would that look like?", False, 0),
        ("", False, 0),
        ("2. RAP and LATS both use MCTS for reasoning. When should an agent use tree search vs simple linear reasoning? What determines the optimal approach?", False, 0),
        ("", False, 0),
        ("3. Valmeekam et al. show LLMs achieve only ~3% on autonomous planning. Does this fundamentally limit agent architectures, or is it resolved by grounding (ReAct)?", False, 0),
        ("", False, 0),
        ("4. The Scaling Agents paper shows MAS can hurt (-70% on PlanCraft). How should practitioners decide between single-agent and multi-agent architectures?", False, 0),
        ("", False, 0),
        ("5. Test-time compute scaling suggests smaller models + search can match larger models. How does this change the economics of deploying agentic systems?", False, 0),
    ],
)

# Reading List
add_text_slide(
    "Reading List: W2 Papers",
    [
        ("Chain-of-Thought & Reasoning", True, 0),
        ("[Wei et al. 22] Chain-of-Thought Prompting, NeurIPS 2022", False, 1),
        ("[Wang et al. 23] Self-Consistency, ICLR 2023", False, 1),
        ("[Zhou et al. 23] Least-to-Most Prompting, ICLR 2023", False, 1),
        ("[Hao et al. 23] RAP, EMNLP 2023", False, 1),
        ("[Snell et al. 25] Test-Time Compute Scaling, ICML 2025", False, 1),
        ("", False, 0),
        ("Search & Planning", True, 0),
        ("[Yao et al. 23] Tree of Thoughts, NeurIPS 2023", False, 1),
        ("[Wang et al. 23] Plan-and-Solve, ACL 2023", False, 1),
        ("[Zhou et al. 24] LATS, ICML 2024", False, 1),
        ("[Besta et al. 24] Graph of Thoughts, AAAI 2024", False, 1),
        ("[Huang et al. 24] Planning Survey, arXiv 2024", False, 1),
        ("[Valmeekam et al. 23] Planning Abilities, NeurIPS 2023", False, 1),
    ],
    bullet_size=Pt(9),
)

add_text_slide(
    "Reading List: W2 Papers (cont.)",
    [
        ("Grounded Acting & ReAct", True, 0),
        ("[Yao et al. 23] ReAct, ICLR 2023", False, 1),
        ("[Shinn et al. 23] Reflexion, NeurIPS 2023", False, 1),
        ("[Chen et al. 23] FireAct, arXiv 2023", False, 1),
        ("[Kim et al. 25] Scaling Agent Systems, arXiv 2025", False, 1),
    ],
    bullet_size=Pt(9),
)

# ─── Save ───
output_path = "/Users/sungju/Agentic AI/slides/W2_Agentic_Architecture_Patterns.pptx"
prs.save(output_path)
print(f"\nSaved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
